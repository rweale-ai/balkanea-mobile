import copy
from pptx import Presentation
from pptx.util import Emu, Pt

PATH = 'Balkanea-Mobile-Sandbox-Readiness-Update.pptx'
prs = Presentation(PATH)
s = prs.slides[5]


def shape_by_name(slide, name):
    matches = [sh for sh in slide.shapes if sh.name == name]
    return matches[0]


def set_single_run_text(shape, new_text):
    tf = shape.text_frame
    first_para = tf.paragraphs[0]
    proto = first_para.runs[0] if first_para.runs else None
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    for r in list(first_para.runs):
        r._r.getparent().remove(r._r)
    run = first_para.add_run()
    run.text = new_text
    if proto is not None:
        run.font.size = proto.font.size
        run.font.bold = proto.font.bold
        run.font.name = proto.font.name
        try:
            if proto.font.color and proto.font.color.type is not None:
                run.font.color.rgb = proto.font.color.rgb
        except Exception:
            pass


# Update subtitle: there IS now a real recurring infra cost, not just fixed fees.
set_single_run_text(shape_by_name(s, 'TextBox 3'),
    'Claude performs all engineering — the cost drivers are calendar time, a small set of fixed fees, and the hotel-DB’s own hosting bill.')

# --- Clone row 1 (Rounded Rectangle 6/7, TextBox 8/9/10) to build a new row for DB hosting ---
row1_names = ['Rounded Rectangle 6', 'Rounded Rectangle 7', 'TextBox 8', 'TextBox 9', 'TextBox 10']
spTree = s.shapes._spTree
clones = []
for name in row1_names:
    src = shape_by_name(s, name)
    new_el = copy.deepcopy(src._element)
    spTree.append(new_el)
    clones.append(new_el)

# Re-fetch the new shapes via python-pptx wrappers (last 5 shapes added)
new_shapes = list(s.shapes)[-5:]
new_container, new_badge_rect, new_badge_text, new_title, new_desc = new_shapes
for sh, suffix in zip(new_shapes, ['Container', 'BadgeRect', 'BadgeText', 'Title', 'Desc']):
    sh.name = f'DBHosting {suffix}'

# --- New 5-row vertical layout: shorter cards, tighter gaps, same left/width ---
ROW_H = Pt(0.84 * 72)  # 0.84in
GAP = Pt(0.10 * 72)
TOP0 = Pt(2.35 * 72)

rows_in_order = [
    ('Rounded Rectangle 6', 'Rounded Rectangle 7', 'TextBox 8', 'TextBox 9', 'TextBox 10'),   # Engineering/AI labor
    (new_container.name, new_badge_rect.name, new_badge_text.name, new_title.name, new_desc.name),  # DB hosting (NEW)
    ('Rounded Rectangle 11', 'Rounded Rectangle 12', 'TextBox 13', 'TextBox 14', 'TextBox 15'),  # Testing & design
    ('Rounded Rectangle 16', 'Rounded Rectangle 17', 'TextBox 18', 'TextBox 19', 'TextBox 20'),  # Store & account fees
    ('Rounded Rectangle 21', 'Rounded Rectangle 22', 'TextBox 23', 'TextBox 24', 'TextBox 25'),  # What's NOT a dollar cost
]

# Original internal offsets (from container top) for badge/title/desc — reused for every row.
BADGE_OFFSET = Pt(0.16 * 72)
TITLE_OFFSET = Pt(0.1997 * 72)
DESC_OFFSET = Pt(0.51 * 72)

for i, (cont_n, badge_n, badgetext_n, title_n, desc_n) in enumerate(rows_in_order):
    top = TOP0 + i * (ROW_H + GAP)
    cont = shape_by_name(s, cont_n)
    cont.top = int(top)
    cont.height = int(ROW_H)

    badge_rect = shape_by_name(s, badge_n)
    badge_text = shape_by_name(s, badgetext_n)
    badge_rect.top = int(top + BADGE_OFFSET)
    badge_text.top = int(top + BADGE_OFFSET)

    title = shape_by_name(s, title_n)
    title.top = int(top + TITLE_OFFSET)

    desc = shape_by_name(s, desc_n)
    desc.top = int(top + DESC_OFFSET)

# --- Set content for the new DB hosting row ---
set_single_run_text(shape_by_name(s, new_title.name), 'Hotel database hosting (Supabase)')
set_single_run_text(shape_by_name(s, new_desc.name),
    'Sandbox running now at ~$25/mo (Pro plan, Medium compute). Production baseline ~$75/mo once live, '
    'plus an optional HA read replica (~$50–60/mo) if traffic justifies it. Backups skipped — hotel data '
    'is fully re-derivable from RateHawk, not transactional.')

badge_shape = shape_by_name(s, new_badge_rect.name)
badge_text_shape = shape_by_name(s, new_badge_text.name)
label = 'LIVE — SCALES WITH USAGE'
right_edge = Emu(badge_shape.left) + Emu(badge_shape.width)
char_w = Emu(Pt(5.3))
new_width = int(char_w) * len(label) + Emu(Pt(14))
new_left = int(right_edge) - new_width
badge_shape.left = new_left
badge_shape.width = new_width
badge_text_shape.left = new_left
badge_text_shape.width = new_width
set_single_run_text(badge_text_shape, label)

prs.save(PATH)
print('Saved', PATH)
