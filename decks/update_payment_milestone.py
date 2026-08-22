"""
One-time content update: reflects the 8/20 milestone with Hristijan
(real end-to-end payment cleared through NLB via the Balkanea Payment
Bridge plugin, from the actual mobile app booking flow) and the one
remaining open item (Bankart's hosted card page isn't mobile-responsive).

Edits text in place; does not touch slide structure (no insert/delete of
slides), so this carries none of the corruption risk of the earlier
slide-generation scripts.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor

PATH = 'Balkanea-Mobile-Sandbox-Readiness-Update.pptx'
prs = Presentation(PATH)

GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_FONT = RGBColor(0x2E, 0x6B, 0x57)


def set_text(shape, new_text):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    proto = p.runs[0]
    proto_size = proto.font.size
    proto_bold = proto.font.bold
    proto_name = proto.font.name
    proto_color = proto.font.color.rgb
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    for extra_p in list(tf.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)
    run = p.add_run()
    run.text = new_text
    run.font.size = proto_size
    run.font.bold = proto_bold
    run.font.name = proto_name
    run.font.color.rgb = proto_color
    return run


# ---------- Slide index 3: PAYMENTS ----------
s = prs.slides[3]
shapes = s.shapes

set_text(shapes[1], 'Bank checkout works end-to-end — confirmed live with NLB')
set_text(
    shapes[2],
    "Confirmed 8/20: a real test-card payment cleared through NLB via Bankart, "
    "from the app's booking flow to a synced Salesforce lead. What's left: "
    "making Bankart's hosted card page render correctly on phones and tablets.",
)

# "BLOCKED ON PARTNER" -> resolved, recolor to green to match CONFIRMED/LIVE
shapes[11].fill.solid()
shapes[11].fill.fore_color.rgb = GREEN_FILL
set_text(shapes[12], 'CONFIRMED LIVE')
shapes[12].text_frame.paragraphs[0].runs[0].font.color.rgb = GREEN_FONT

set_text(shapes[13], 'Balkanea Payment Bridge — built and confirmed working end-to-end')
set_text(
    shapes[14],
    "Hristijan's ‘Balkanea Payment Bridge’ plugin (signed HMAC links, Bankart "
    "Gateway V3 hosted card form, async Notify webhook) is live. Confirmed 8/20: a "
    "real test-card payment cleared through NLB, initiated from the actual mobile "
    "app booking flow — 3DS, async postback, booking auto-confirm, and Salesforce "
    "sync all verified working end-to-end.",
)

set_text(shapes[18], 'Outstanding: mobile-responsive card page')
set_text(
    shapes[19],
    "Bankart's hosted card page isn't mobile-responsive today — input fields "
    "misalign on phone and tablet, iOS and Android. That's the one open item "
    "blocking a clean guest-facing checkout; on Hristijan's side to fix.",
)

# ---------- Slide index 7: LIVE DEMO ----------
s = prs.slides[7]
shapes = s.shapes

set_text(shapes[1], 'End-to-end payment now works — from the app')
set_text(
    shapes[2],
    "Confirmed 8/20: a real test card cleared through NLB, from the mobile app's "
    "own booking flow to a synced Salesforce lead.",
)

set_text(shapes[6], 'What happened')
set_text(
    shapes[7],
    "A real booking in the app → signed Bankart payment link → Bankart's "
    "hosted card form → NLB processing with 3DS → async Notify webhook confirms "
    "the booking → Salesforce lead synced. The actual app flow, start to finish "
    "— not a browser demo.",
)

set_text(shapes[10], 'What it proves: the full stack is real')
set_text(
    shapes[11],
    "This wasn't a sandboxed demo — a real card ran through NLB via Bankart, "
    "initiated from the actual mobile app, ending in a confirmed booking and a "
    "synced CRM lead. Every layer (signed link, hosted card form, async webhook, "
    "Supabase state) is now proven live together, not just individually tested.",
)

set_text(shapes[14], "What's still open: the card page isn't mobile-responsive")
set_text(
    shapes[15],
    "Bankart's hosted card page (Hristijan's side, outside our control) has "
    "alignment issues on phone and tablet screens — the one open item before "
    "this is a clean guest-facing checkout on iOS and Android.",
)
set_text(shapes[16], 'OPEN ITEM')

# Orphan "ACTION" badge (shape 17) has no accompanying card/heading/body and
# predates this edit -- remove rather than leave dangling.
orphan = shapes[17]
orphan._element.getparent().remove(orphan._element)

# ---------- Slide index 12: STILL OPEN ----------
s = prs.slides[12]
shapes = s.shapes

set_text(
    shapes[13],
    "Confirmed 8/20: a real end-to-end payment cleared through NLB via Bankart, "
    "from the app's booking flow to a synced Salesforce lead. Only remaining item: "
    "Bankart's hosted card page isn't mobile-responsive on phone/tablet — "
    "Hristijan-side fix. RateHawk's booking flow draws on its own deposit balance, "
    "not the guest's card, so RateHawk build work doesn't wait on this.",
)
set_text(
    shapes[16],
    "Email to RateHawk support on key 954 status; mobile-responsive fix timeline "
    "for Bankart's hosted card page, from Hristijan",
)

prs.save(PATH)
print('Saved.')
