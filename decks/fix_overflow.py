from pptx import Presentation

PATH = 'Balkanea-Mobile-Sandbox-Readiness-Update.pptx'
prs = Presentation(PATH)


def shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError(name)


def set_single_run_text(shape, new_text):
    tf = shape.text_frame
    first_para = tf.paragraphs[0]
    proto = first_para.runs[0] if first_para.runs else None
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    for r in list(first_para.runs):
        r._r.getparent().remove(r._r)
    run = first_para.add_run()
    run.text = new_text
    if proto is not None:
        run.font.size = proto.font.size
        run.font.bold = proto.font.bold
        run.font.name = proto.font.name
        try:
            if proto.font.color and proto.font.color.type is not None:
                run.font.color.rgb = proto.font.color.rgb
        except Exception:
            pass


# Slide 4 (index 3), card 1 desc — was overflowing the card border
s = prs.slides[3]
new_text = ("Hristijan built and documented the ‘Mobile Payment Bridge’ plugin — WebView opens a "
            "guest WooCommerce order-pay page, status polled via REST. Verified live against a real staging "
            "order. But it's built entirely on WooCommerce's order-key model, which ‘no WooCommerce in "
            "the architecture’ now rules out.")
print('slide4 card1 len:', len(new_text))
set_single_run_text(shape_by_name(s, 'TextBox 15'), new_text)

# Slide 8 (index 7), footer line — was wrapping into the page-number/branding footer
s = prs.slides[7]
footer = shape_by_name(s, 'TextBox 18')
p = footer.text_frame.paragraphs[0]
runs = list(p.runs)
new_footer = ("no Google Play account yet · Apple review ~3–5 days · Google's 14-day closed-testing "
              "minimum · ~1 month+ total, RateHawk cert now runs in parallel.")
print('slide8 footer len:', len(new_footer))
runs[1].text = new_footer

prs.save(PATH)
print('Saved', PATH)
