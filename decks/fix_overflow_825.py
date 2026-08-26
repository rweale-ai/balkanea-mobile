# -*- coding: utf-8 -*-
"""
Fixes found on render-verification pass:

1. Slide 2 (Search & Booking): the first edit pass wrote the "booking flow proven"
   content into shapes[16..19], which is actually CARD 3 (the "hotel database is
   live" card) -- clobbering that content. The real target was CARD 4
   (shapes[21..24], badge "Open" -> the stale "Search-flow now settled..." card).
   Restores card 3, and puts a shortened "booking flow proven" edit into card 4.
2. Slide 3 (Payments): subtitle overflowed into the first card -- shortened to 2 lines.
3. Slide 6 (Live Demo... no -- Alignment slide, index 5): card 2 body was flush
   against the IP caption line -- trimmed slightly for margin.
4. Slide 7 (Planning): Phase 1 and Phase 2 card bodies overflowed their cards --
   shortened both to fit within 3 lines like Phase 3/4 do.
5. Slide 9 (Cost 1): "Hotel database hosting" card body overflowed -- shortened.
6. Slide 11 (Still Open... now index 12): final "Store submission" line wrapped to
   2 lines and collided with the footer -- shortened back to 1 line.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)

GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_FONT = RGBColor(0x2E, 0x6B, 0x57)


def set_text(shape, new_text):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    proto = p.runs[0]
    proto_size = proto.font.size
    proto_bold = proto.font.bold
    proto_name = proto.font.name
    proto_color = proto.font.color.rgb
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    for extra_p in list(tf.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)
    run = p.add_run()
    run.text = new_text
    run.font.size = proto_size
    run.font.bold = proto_bold
    run.font.name = proto_name
    run.font.color.rgb = proto_color
    return run


def recolor_badge(fill_shape, label_shape, text, positive):
    fill_shape.fill.solid()
    fill_shape.fill.fore_color.rgb = GREEN_FILL if positive else fill_shape.fill.fore_color.rgb
    set_text(label_shape, text)
    label_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = GREEN_FONT if positive else label_shape.text_frame.paragraphs[0].runs[0].font.color.rgb


# =====================================================================
# FIX 1 -- Slide 2 (Search & Booking): restore card 3, fix card 4
# =====================================================================
s = prs.slides[2]
sh = s.shapes

# Restore card 3 (shapes[16..19]) to its original "hotel DB is live" content.
recolor_badge(sh[16], sh[17], "RESOLVED", True)
set_text(sh[18], "Mobile's own hotel database is live")
set_text(
    sh[19],
    "850,218 hotels / 3,930,128 rooms imported and indexed across all 10 target "
    "countries (0 parse errors). Region-filtered, sorted queries benchmarked at "
    "~0.2ms server-side. Hristijan already has Developer access to query it directly.",
)

# Card 4 (shapes[21..24]) -- badge fill is sh[21], label is sh[22]; title sh[23], body sh[24].
recolor_badge(sh[21], sh[22], "CONFIRMED", True)
set_text(sh[23], "Full booking flow proven end-to-end in sandbox, including multi-room")
set_text(
    sh[24],
    "Search \u2192 prebook \u2192 create order \u2192 finish now runs successfully against "
    "RateHawk's certification test hotels (Conrad LA, Rosa Bell Motel) -- all 7 "
    "mandatory scenarios passed, real order IDs captured. Live-app wiring and the "
    "room-by-room guest picker UI are next -- see Risks and Plan.",
)

# =====================================================================
# FIX 2 -- Slide 3 (Payments): shorten subtitle to 2 lines
# =====================================================================
s = prs.slides[3]
sh = s.shapes
set_text(
    sh[2],
    "Confirmed 8/20: a real test-card payment cleared through NLB via Bankart, from "
    "the app's booking flow to a synced Salesforce lead. Settles in MKD only -- see "
    "Risks for the one open item.",
)

# =====================================================================
# FIX 3 -- Slide 5 (Alignment & Recommendation): trim card 2 body slightly
# =====================================================================
s = prs.slides[5]
sh = s.shapes
set_text(
    sh[11],
    "For the actual requirement -- a static IP RateHawk's production certification "
    "will accept -- the AWS build is live, verified, and ~$5\u20138/mo today, with no "
    "further engineering needed. Ivan's fuller proposal suits a mature, high-traffic "
    "system, but adopting it as specified means more cost and an unresolved question "
    "of who runs it -- neither of which certification itself requires. Fastest path: "
    "use what's already built and verified.",
)

# =====================================================================
# FIX 4 -- Slide 7 (Planning): shorten Phase 1 and Phase 2 bodies
# =====================================================================
s = prs.slides[7]
sh = s.shapes
set_text(
    sh[9],
    "Answer RateHawk's Pre-Certification Checklist in full (payment type, RPM limits, "
    "static-data sync, price-parsing field, booking-success signal, full error-"
    "handling matrix), submit the AWS static IP for whitelisting, and hand back the "
    "real order IDs from the 7 passed test scenarios. RateHawk's own review time after "
    "that isn't ours to set.",
)
set_text(
    sh[14],
    "Replace lib/ratehawk.ts's simulated stub with the real prebook \u2192 charge \u2192 "
    "finish sequence, reconciling that RateHawk charges Balkanea's deposit balance "
    "while Bankart charges the guest's card. Build the room-by-room guest configurator "
    "that search and booking currently lack entirely (multi-file, not small). Apple's "
    "OAuth secret + Google Play account (Ray-owned) proceed in parallel.",
)

# =====================================================================
# FIX 5 -- Slide 9 (Cost 1): shorten hotel DB hosting card body
# =====================================================================
s = prs.slides[9]
sh = s.shapes
set_text(
    sh[29],
    "Sandbox running now at ~$25/mo (Pro plan). Production baseline ~$75/mo once "
    "live. Plus the AWS relay RateHawk's production certification requires: ~$8/mo "
    "base, ~$22/mo with HA -- both now in the cost model.",
)

# =====================================================================
# FIX 6 -- Slide 12 (Still Open): shorten final store-submission line to 1 line
# =====================================================================
s = prs.slides[12]
sh = s.shapes
set_text(
    sh[17],
    "Store submission, once ready: no Google Play account yet \u00b7 Apple review "
    "~3\u20135 days \u00b7 Google's 14-day closed-testing minimum \u00b7 RateHawk cert "
    "runs in parallel.",
)

prs.save(PATH)
print("Fixes saved.")
