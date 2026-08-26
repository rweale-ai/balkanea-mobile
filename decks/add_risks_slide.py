# -*- coding: utf-8 -*-
"""
Adds a new RISKS slide, inserted after the two Cost slides (index 10) and before
"Still Open" (currently index 11) -- so it lands at index 11, total slides 12 -> 13.

Covers the four gaps between what this update treats as working (per Ray's
direction) and what's actually verified today:
  1. Multi-room booking -- zero real product support in the app (the biggest gap).
  2. Real-hotel payments -- intermittent settlement failure, unresolved (error_code 1000).
  3. RateHawk's real booking flow -- proven only in standalone scripts, not wired into the app.
  4. External review gates -- RateHawk certification + store submission timelines.

Reuses the card/badge/textbox pattern from add_ip_whitelisting_slide.py.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)

INSERT_AT = 11  # after the 2nd Cost slide (index 10), before "Still Open" (currently index 11)


def drop_slide_properly(prs, index):
    xml_slides = prs.slides._sldIdLst
    slide_id_el = list(xml_slides)[index]
    rId = slide_id_el.get(qn('r:id'))
    prs.part.drop_rel(rId)
    xml_slides.remove(slide_id_el)


# Re-run safety: drop a previously-added copy of this slide first.
for i, s in enumerate(prs.slides):
    for shp in s.shapes:
        if shp.has_text_frame and 'RISKS' in shp.text_frame.text and '\u2014' in shp.text_frame.text:
            drop_slide_properly(prs, i)
            break

CREAM = RGBColor(0xFA, 0xF7, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE7, 0xE0, 0xD3)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x65, 0x58)
FOOTER_GRAY = RGBColor(0x9B, 0x93, 0x84)
BROWN = RGBColor(0x7A, 0x3E, 0x0F)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD9)
AMBER_TEXT = RGBColor(0xB8, 0x79, 0x0A)
GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_TEXT = RGBColor(0x2E, 0x6B, 0x57)

layout = prs.slides[3].slide_layout  # same blank layout used for other content slides
new_slide = prs.slides.add_slide(layout)
new_slide.background.fill.solid()
new_slide.background.fill.fore_color.rgb = CREAM
sp = new_slide.shapes


def textbox(left, top, width, height, text, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = sp.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.name = 'Segoe UI Semibold' if bold else 'Segoe UI'
    run.font.color.rgb = color
    return box


def badge(rect_ref, label, positive=True):
    w = Emu(int(Emu(Pt(5.3)) * len(label) + Emu(Pt(16))))
    h = Inches(0.20)
    left = Emu(int(Emu(rect_ref.left) + Emu(rect_ref.width) - w - Pt(8)))
    top = Emu(int(Emu(rect_ref.top) + Pt(7)))
    b = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    b.adjustments[0] = 0.5
    b.fill.solid()
    b.fill.fore_color.rgb = GREEN_FILL if positive else AMBER_FILL
    b.line.fill.background()
    b.shadow.inherit = False
    tf = b.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(7.5)
    run.font.bold = True
    run.font.name = 'Segoe UI Semibold'
    run.font.color.rgb = GREEN_TEXT if positive else AMBER_TEXT
    return b


def card(left, top, width, height, title, desc, badge_label, badge_positive):
    rect = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.09
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = CARD_BORDER
    rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    pad = Pt(10)
    title_top = Emu(int(top + Inches(0.12)))
    textbox(Emu(int(left + pad)), title_top, Emu(int(width - 2 * pad - Inches(1.6))), Inches(0.22), title, Pt(11.5), bold=True, color=NAVY)
    desc_top = Emu(int(title_top + Inches(0.28)))
    textbox(Emu(int(left + pad)), desc_top, Emu(int(width - 2 * pad)), Emu(int(top + height - desc_top - Pt(6))), desc, Pt(9.5), color=GRAY)
    if badge_label:
        badge(rect, badge_label, positive=badge_positive)
    return rect


# ---------------------------------------------------------------------------
# Header
# ---------------------------------------------------------------------------
textbox(Inches(0.55), Inches(0.50), Inches(12.23), Inches(0.35), '13  \u2014  RISKS', Pt(12), bold=True, color=BROWN)
textbox(Inches(0.55), Inches(0.82), Inches(12.23), Inches(0.70),
        'What "treated as working" assumes, and where it could break', Pt(27), bold=True, color=NAVY)
textbox(Inches(0.55), Inches(1.55), Inches(11.81), Inches(0.55),
        'Payments and RateHawk booking are shown as working elsewhere in this deck for planning purposes. '
        'This is the honest gap between that and what is actually verified today.',
        Pt(13.5), color=GRAY)
textbox(Inches(0.55), Inches(7.14), Inches(8.00), Inches(0.30), 'BALKANEA MOBILE  \u2014  SANDBOX / TEST READINESS', Pt(9), color=FOOTER_GRAY)
textbox(Inches(11.78), Inches(7.14), Inches(1.00), Inches(0.30), '13', Pt(9), color=FOOTER_GRAY, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# Cards
# ---------------------------------------------------------------------------
CARD_X, CARD_W = Inches(0.55), Inches(12.23)
TOP0 = Inches(2.35)
CARD_H = Inches(1.02)
GAP = Inches(0.15)

card(CARD_X, TOP0, CARD_W, CARD_H,
    'Multi-room booking has no real product support in the app today',
    "RateHawk's booking API proves multi-room works end-to-end (real order IDs, mixed child ages tested) "
    "-- but the app itself has no room-by-room guest picker, a single flat guest-name field, and rooms "
    "passed through as an unused integer. Treated as working elsewhere in this deck per Ray's direction; "
    "it is not built. Scoped as its own multi-file build item on the Plan slide, not a small addition.",
    'BIGGEST GAP', False)

card(CARD_X, Emu(int(TOP0 + CARD_H + GAP)), CARD_W, CARD_H,
    'Real-hotel payments still fail intermittently, unresolved',
    "Bankart's merchant account is confirmed MKD-only. A real-hotel MKD payment attempt on 8/25 still "
    "failed with a generic Bankart-side error (error_code 1000, uuid: null) after currency, reference "
    "length, and format were all ruled out as causes. Root cause sits in Hristijan's plugin or Bankart's "
    "gateway, not the app -- treated as working for this update, but not yet clean on a real hotel.",
    'UNRESOLVED', False)

card(CARD_X, Emu(int(TOP0 + 2 * (CARD_H + GAP))), CARD_W, CARD_H,
    "RateHawk's real booking flow isn't wired into the app yet",
    "Search \u2192 prebook \u2192 finish is proven only in standalone test scripts against RateHawk's "
    "sandbox -- lib/ratehawk.ts still runs the original simulated stub. Wiring it in has a real failure "
    "mode to design around first: RateHawk charges Balkanea's own deposit balance, not the guest's card, "
    "so a guest could be charged by Bankart while RateHawk's own booking finish independently fails.",
    'INTEGRATION GAP', False)

card(CARD_X, Emu(int(TOP0 + 3 * (CARD_H + GAP))), CARD_W, CARD_H,
    "External review gates: RateHawk certification and store submission",
    "RateHawk won't issue production keys until Balkanea answers its full certification questionnaire "
    "(payment type, RPM limits, static-data sync method, a complete booking error-handling matrix) and "
    "the AWS static IP is whitelisted -- review time after that is RateHawk's, not ours. Google Play's "
    "14-day closed-testing clock can't start until the Play Developer account exists (still not created).",
    'RATEHAWK & STORE-PACED', False)

prs.save(PATH)

# ---------------------------------------------------------------------------
# Move the new slide (currently last) into position INSERT_AT
# ---------------------------------------------------------------------------
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
new_el = slides[-1]
xml_slides.remove(new_el)
xml_slides.insert(INSERT_AT, new_el)
prs.save(PATH)
print('Inserted RISKS slide at index', INSERT_AT, '- total slides now:', len(prs.slides))
