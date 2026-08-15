"""Update Balkanea-Mobile-Sandbox-Readiness.pptx with latest info, same structure."""
import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

SRC = 'Balkanea-Mobile-Sandbox-Readiness.pptx'
OUT = 'Balkanea-Mobile-Sandbox-Readiness-Update.pptx'

AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD9)
AMBER_TEXT = RGBColor(0xB8, 0x79, 0x0A)
GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_TEXT = RGBColor(0x2E, 0x6B, 0x57)

prs = Presentation(SRC)
slides = prs.slides


def set_single_run_text(shape, new_text):
    """Replace text of a (single-run-per-paragraph) text box, keeping the first run's formatting."""
    tf = shape.text_frame
    first_para = tf.paragraphs[0]
    proto = first_para.runs[0] if first_para.runs else None
    # clear all paragraphs except the first
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    # clear runs in first paragraph, keep pPr
    for r in list(first_para.runs):
        r._r.getparent().remove(r._r)
    run = first_para.add_run()
    run.text = new_text
    if proto is not None:
        run.font.size = proto.font.size
        run.font.bold = proto.font.bold
        run.font.name = proto.font.name
        try:
            if proto.font.color and proto.font.color.type is not None:
                run.font.color.rgb = proto.font.color.rgb
        except Exception:
            pass


def set_multiline_text(shape, lines):
    """Replace text of a text box with N paragraphs of plain text, cloning the first paragraph's pPr/run formatting for every line."""
    tf = shape.text_frame
    template_p = tf.paragraphs[0]._p
    proto_run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else None
    txBody = tf._txBody
    # remove all existing <a:p>
    for p in list(tf.paragraphs):
        txBody.remove(p._p)
    for line in lines:
        new_p = copy.deepcopy(template_p)
        # strip existing runs from the clone
        for r in list(new_p.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r')):
            new_p.remove(r)
        txBody.append(new_p)
        from pptx.text.text import _Paragraph
        para = _Paragraph(new_p, tf)
        run = para.add_run()
        run.text = line
        if proto_run is not None:
            run.font.size = proto_run.font.size
            run.font.bold = proto_run.font.bold
            run.font.name = proto_run.font.name
            try:
                if proto_run.font.color and proto_run.font.color.type is not None:
                    run.font.color.rgb = proto_run.font.color.rgb
            except Exception:
                pass


def set_badge(rect_shape, text_shape, label, positive):
    fill = GREEN_FILL if positive else AMBER_FILL
    text_color = GREEN_TEXT if positive else AMBER_TEXT
    rect_shape.fill.solid()
    rect_shape.fill.fore_color.rgb = fill
    # resize badge box to roughly fit new label, keep right edge fixed
    right_edge = Emu(rect_shape.left) + Emu(rect_shape.width)
    char_w = Emu(Pt(5.3))  # ~7.5pt bold uppercase avg width per char
    new_width = int(char_w) * len(label) + Emu(Pt(14))  # padding
    new_left = int(right_edge) - new_width
    rect_shape.left = new_left
    rect_shape.width = new_width
    text_shape.left = new_left
    text_shape.width = new_width
    set_single_run_text(text_shape, label)
    for p in text_shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = text_color


def shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError(name)


# ---------- SLIDE 1 (index 0): TITLE ----------
s = slides[0]
set_single_run_text(shape_by_name(s, 'TextBox 1'), 'BALKANEA MOBILE  —  SANDBOX / TEST READINESS UPDATE')
set_multiline_text(shape_by_name(s, 'TextBox 2'), [
    "What's changed since the last review",
    'and what still needs a decision',
])
set_multiline_text(shape_by_name(s, 'TextBox 3'), [
    'Hotel database is live, RateHawk sandbox is being provisioned, and payments hit a',
    'new architecture call — prepared for review with Dragana & Hristijan',
])

# meta line: multi-run paragraph (label bold / value regular, repeated x3)
meta = shape_by_name(s, 'TextBox 4')
tf = meta.text_frame
p = tf.paragraphs[0]
proto_runs = list(p.runs)
label_font = proto_runs[0].font
value_font = proto_runs[1].font
new_meta = [
    ('Stage:  ', True), ('TestFlight (v1.0.7, build 11)      ', False),
    ('Systems:  ', True), ('RateHawk · Bank checkout (architecture TBD) · Supabase auth      ', False),
    ('Prepared:  ', True), ('August 2026 (update)', False),
]
for r in list(p.runs):
    r._r.getparent().remove(r._r)
for text, bold in new_meta:
    run = p.add_run()
    run.text = text
    src_font = label_font if bold else value_font
    run.font.size = src_font.size
    run.font.bold = src_font.bold
    run.font.name = src_font.name
    run.font.color.rgb = src_font.color.rgb

# ---------- SLIDE 2 (index 1): 01 AUTHENTICATION ----------
s = slides[1]
set_single_run_text(shape_by_name(s, 'TextBox 2'), "Google sign-in is live; Apple is the remaining gap")
set_single_run_text(shape_by_name(s, 'TextBox 3'),
    "Google OAuth is confirmed working end to end. What's left is Apple's provider secret and the Google Play account.")

set_badge(shape_by_name(s, 'Rounded Rectangle 7'), shape_by_name(s, 'TextBox 8'), 'RESOLVED', True)
set_single_run_text(shape_by_name(s, 'TextBox 9'), 'Google OAuth is live')
set_single_run_text(shape_by_name(s, 'TextBox 10'),
    'Confirmed via the Supabase Management API: real client credentials are configured and the balkanea://auth/callback redirect is already allow-listed. No action needed.')

set_badge(shape_by_name(s, 'Rounded Rectangle 12'), shape_by_name(s, 'TextBox 13'), 'BLOCKED', False)
set_single_run_text(shape_by_name(s, 'TextBox 14'), 'Native Sign-in with Apple — built, not switched on')
set_single_run_text(shape_by_name(s, 'TextBox 15'),
    "Native code shipped 8/5, but Supabase's Apple provider has no client secret yet — it'll fail at runtime until that's set. Needs a Services ID + signed JWT from Ray's own Apple Developer account; Claude can't do this part.")

# Card 3 (Google Play Developer account) unchanged — leave as-is.

# ---------- SLIDE 3 (index 2): 02 SEARCH & BOOKING ----------
s = slides[2]
set_single_run_text(shape_by_name(s, 'TextBox 2'), 'Hotel database is live; RateHawk sandbox is provisioning')
set_single_run_text(shape_by_name(s, 'TextBox 3'),
    "The hotel-DB half of this is done. What's left: RateHawk's sandbox key, the search-flow decision, and the booking endpoints.")

set_badge(shape_by_name(s, 'Rounded Rectangle 7'), shape_by_name(s, 'TextBox 8'), 'IN PROGRESS', True)
set_single_run_text(shape_by_name(s, 'TextBox 9'), 'RateHawk sandbox — separate contract confirmed, key incoming')
set_single_run_text(shape_by_name(s, 'TextBox 10'),
    "Ivan's proposed contract slug (balkanea-mobile) was accepted by RateHawk today (8/14) — they're creating the contract now and handing off to their API support team to issue the sandbox key. New wrinkle: RateHawk's API Launch team also requires a certification step (demo + checklist) before issuing a production key — a gate not in the original plan.")

set_badge(shape_by_name(s, 'Rounded Rectangle 12'), shape_by_name(s, 'TextBox 13'), 'CONFIRMED', True)
set_single_run_text(shape_by_name(s, 'TextBox 14'), 'B2C vs. B2B is a pricing/availability split, not a separate hotel file')
set_single_run_text(shape_by_name(s, 'TextBox 15'),
    "The master dataset (~4M hotels) is shared between B2C and B2B — RateHawk exposes only price/availability differently per credential type. The mobile app's code comment describing a separate 'B2C-approved hotel file' still needs a one-line update.")

set_badge(shape_by_name(s, 'Rounded Rectangle 17'), shape_by_name(s, 'TextBox 18'), 'RESOLVED', True)
set_single_run_text(shape_by_name(s, 'TextBox 19'), "Mobile's own hotel database is live")
set_single_run_text(shape_by_name(s, 'TextBox 20'),
    '850,218 hotels / 3,930,128 rooms imported and indexed across all 10 target countries (0 parse errors). Region-filtered, sorted queries benchmarked at ~0.2ms server-side — round trip is now ~215ms of pure network latency, not DB work. Hristijan already has Developer access to query it directly.')

set_single_run_text(shape_by_name(s, 'TextBox 24'), 'Search-flow design, and booking endpoints still to build')
set_single_run_text(shape_by_name(s, 'TextBox 25'),
    "Still undecided: web's live-search-then-enrich pattern vs. the DB-first shortcut (query our own DB first, no price shown until a hotel is opened). Backend also needs RateHawk's mobile-specific SERP/Prebook/Webhook workflow, not just web's lib/ratehawk.js pattern with swapped credentials — none of the room-rate, pre-book, or create-order calls exist yet.")

# ---------- SLIDE 4 (index 3): 03 PAYMENTS ----------
s = slides[3]
set_single_run_text(shape_by_name(s, 'TextBox 2'), 'Bank checkout — the WooCommerce path is ruled out')
set_single_run_text(shape_by_name(s, 'TextBox 3'),
    'A working WebView integration got built and verified — then a hard constraint ruled it out. Now waiting on the real alternative.')

set_badge(shape_by_name(s, 'Rounded Rectangle 12'), shape_by_name(s, 'TextBox 13'), 'SUPERSEDED', False)
set_single_run_text(shape_by_name(s, 'TextBox 14'), 'WooCommerce bridge — built, verified, ruled out')
set_single_run_text(shape_by_name(s, 'TextBox 15'),
    "Hristijan built and documented the ‘Mobile Payment Bridge’ plugin — WebView opens a guest WooCommerce order-pay page, status polled via REST. Verified live against a real staging order (Bankart, sim mode). But it's built entirely on WooCommerce's order-key model — and ‘no WooCommerce in the mobile architecture’ is now a hard constraint, so this won't be the final integration despite working.")

set_badge(shape_by_name(s, 'Rounded Rectangle 17'), shape_by_name(s, 'TextBox 18'), 'BLOCKED ON PARTNER', False)
set_single_run_text(shape_by_name(s, 'TextBox 19'), 'NLB-direct — the only path that fits, not yet built')
set_single_run_text(shape_by_name(s, 'TextBox 20'),
    "Hristijan flagged an NLB EPOS route that skips WooCommerce order creation entirely — surfaced in an 8/14 meeting with Ivan, but 'needs to be developed.' No doc or timeline exists yet on either side; it may be what 'Balkanea Payment Links' (bpl-nlb-settings) already does, but that page is access-blocked for Hristijan's account.")

set_badge(shape_by_name(s, 'Rounded Rectangle 22'), shape_by_name(s, 'TextBox 23'), 'ACTION THIS WEEK', False)
set_single_run_text(shape_by_name(s, 'TextBox 24'), 'Asked for the real spec — this week')
set_single_run_text(shape_by_name(s, 'TextBox 25'),
    "Message out to Hristijan (looping Ivan) for the actual NLB documentation, what 'needs to be developed' means concretely, and the key question: does the order-free flow still create a WooCommerce order behind the scenes for reconciliation? That answer decides if the constraint is even satisfiable through this vendor.")

# ---------- SLIDE 5 (index 4): 04 PLANNING ----------
s = slides[4]
set_single_run_text(shape_by_name(s, 'TextBox 3'),
    'Claude does all technical implementation; testing and store-asset design draw on existing team time. Payments and RateHawk are both gated by partners right now, not effort.')

set_single_run_text(shape_by_name(s, 'TextBox 10'),
    "NLB-direct payment spec from Hristijan/Ivan (the new critical path — nothing on payments builds further until it lands), RateHawk sandbox key (contract confirmed 8/14, issuance in progress), and Google Play Developer account creation (still not started).")

set_single_run_text(shape_by_name(s, 'TextBox 15'),
    "Payment gateway built against the NLB-direct spec once it lands (net-new — two earlier attempts were built on wrong assumptions and discarded), RateHawk's mobile-specific room-rate/pre-book/create-order endpoints plus their certification checklist, search-flow wired up, and Apple's OAuth secret configured for the native sign-in code already shipped.")

# Phase 3 (TextBox 19/20) unchanged.

set_single_run_text(shape_by_name(s, 'TextBox 24'), '~1 month+, both stores + RateHawk cert  ·  estimate')
set_single_run_text(shape_by_name(s, 'TextBox 25'),
    "Apple review ~3–5 days (real travel payments = higher scrutiny). Google Play's 14-day mandatory closed-testing period is the long pole. New: RateHawk's own certification (pre-cert checklist + demo access) also gates production API access, running in parallel. Budget at least one rejection-and-resubmit round on each store.")

# ---------- SLIDES 6 & 7 (index 5, 6): COST — unchanged, no new financial info this session ----------

# ---------- SLIDE 8 (index 7): 06 STILL OPEN ----------
s = slides[7]
set_single_run_text(shape_by_name(s, 'TextBox 1'), '06  —  STILL OPEN')
set_single_run_text(shape_by_name(s, 'TextBox 2'), "What's still open, and what's moving")
# TextBox 3 (subtitle) unchanged — the scheduling constraint still stands.

set_single_run_text(shape_by_name(s, 'TextBox 8'),
    'Contract slug confirmed (balkanea-mobile, accepted 8/14) — RateHawk is creating the contract now and handing off to their API support team to issue the sandbox key. Next: confirm the key lands, then start the certification checklist.')

# Row 2 (Search-flow design) unchanged.

set_single_run_text(shape_by_name(s, 'TextBox 14'),
    "Not a WebView-to-existing-checkout question anymore — get the real NLB-direct spec from Hristijan/Ivan and confirm it doesn't require a WooCommerce order behind the scenes. Nothing else on payments moves until that answer lands.")

set_single_run_text(shape_by_name(s, 'TextBox 17'),
    "The real NLB documentation from Hristijan/Ivan, and Ray's own two items — the Google Play Developer account and the Apple OAuth client secret (Services ID + signed JWT) the already-shipped native Sign-in-with-Apple code needs to actually work.")

footer = shape_by_name(s, 'TextBox 18')
p = footer.text_frame.paragraphs[0]
runs = list(p.runs)
runs[0].text = 'Store submission, once ready: '
runs[1].text = ("no Google Play account yet · Apple review ~3–5 days · Google's 14-day closed-testing minimum · "
                "RateHawk certification runs in parallel · ~1 month+ total, revised up from Hristijan's original estimate.")

prs.save(OUT)
print('Saved', OUT)
