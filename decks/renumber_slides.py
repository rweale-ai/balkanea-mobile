import re
from pptx import Presentation
from pptx.util import Inches

PATH = 'Balkanea-Mobile-Sandbox-Readiness-Update.pptx'
prs = Presentation(PATH)

HEADER_RE = re.compile(r'^(\d{2})(\s{2}—\s{2}.+)$')

# slides from index 5 onward (everything after the new IP Whitelisting slide
# at index 4) need their header number and footer page number bumped.
for idx in range(5, len(prs.slides)):
    s = prs.slides[idx]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()

        # Header line, e.g. "04  --  LIVE DEMO" -> renumber to slide's new position
        m = HEADER_RE.match(t)
        if m and len(t) < 40:
            new_num = f'{idx:02d}'
            new_text = new_num + m.group(2)
            p = sh.text_frame.paragraphs[0]
            proto = p.runs[0]
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            run = p.add_run()
            run.text = new_text
            run.font.size = proto.font.size
            run.font.bold = proto.font.bold
            run.font.name = proto.font.name
            run.font.color.rgb = proto.font.color.rgb
            print(f'slide {idx}: header -> {new_text!r}')
            continue

        # Footer page number: small textbox in the bottom-right corner only.
        if sh.left is not None and sh.top is not None:
            if sh.left > Inches(11) and sh.top > Inches(7) and t.isdigit():
                new_text = str(idx + 1)
                p = sh.text_frame.paragraphs[0]
                proto = p.runs[0]
                for r in list(p.runs):
                    r._r.getparent().remove(r._r)
                run = p.add_run()
                run.text = new_text
                run.font.size = proto.font.size
                run.font.bold = proto.font.bold
                run.font.name = proto.font.name
                run.font.color.rgb = proto.font.color.rgb
                print(f'slide {idx}: footer page -> {new_text!r}')

prs.save(PATH)
print('Saved.')
