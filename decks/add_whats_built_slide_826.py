# -*- coding: utf-8 -*-
"""
Adds a "What's Been Built" slide -- full feature/scope summary + an
illustrative traditional-agency cost comparison -- inserted right after the
Implementation Costs slide (index 9, "09 -- COST", which carries the $1,000
engineering-cost figure) so the contrast lands immediately. Pushes the
running-cost/margin Cost slide, Risks, and Still Open each back by one.

Reuses the card/badge/textbox pattern from add_risks_slide.py for visual
consistency; the cost breakdown uses a native pptx table styled to match
Slide 5's (white rows, thin borders, navy header).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)

INSERT_AT = 10  # right after index 9 ("09 -- COST" / implementation costs)


def drop_slide_properly(prs, index):
    xml_slides = prs.slides._sldIdLst
    slide_id_el = list(xml_slides)[index]
    rId = slide_id_el.get(qn('r:id'))
    prs.part.drop_rel(rId)
    xml_slides.remove(slide_id_el)


# Re-run safety: drop a previously-added copy of this slide first.
for i, s in enumerate(prs.slides):
    for shp in s.shapes:
        if shp.has_text_frame and "WHAT'S BEEN BUILT" in shp.text_frame.text.upper():
            drop_slide_properly(prs, i)
            break

CREAM = RGBColor(0xFA, 0xF7, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE7, 0xE0, 0xD3)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x65, 0x58)
FOOTER_GRAY = RGBColor(0x9B, 0x93, 0x84)
BROWN = RGBColor(0x7A, 0x3E, 0x0F)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD9)
AMBER_TEXT = RGBColor(0xB8, 0x79, 0x0A)
GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_TEXT = RGBColor(0x2E, 0x6B, 0x57)

layout = prs.slides[3].slide_layout
new_slide = prs.slides.add_slide(layout)
new_slide.background.fill.solid()
new_slide.background.fill.fore_color.rgb = CREAM
sp = new_slide.shapes


def textbox(left, top, width, height, text, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT, italic=False):
    box = sp.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = 'Segoe UI Semibold' if bold else 'Segoe UI'
    run.font.color.rgb = color
    return box


def multiline_box(left, top, width, height, lines, size, color=GRAY, line_gap=Pt(3)):
    """lines: list of (text, bold) tuples, one per paragraph, bullet-prefixed by caller."""
    box = sp.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = line_gap
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.name = 'Segoe UI Semibold' if bold else 'Segoe UI'
        run.font.color.rgb = NAVY if bold else color
    return box


# ---------------------------------------------------------------------------
# Header
# ---------------------------------------------------------------------------
textbox(Inches(0.55), Inches(0.50), Inches(12.23), Inches(0.35), '09  \u2014  WHAT\u2019S BEEN BUILT', Pt(12), bold=True, color=BROWN)
textbox(Inches(0.55), Inches(0.82), Inches(12.23), Inches(0.55),
        'Full scope, real codebase size, and what it would have cost elsewhere', Pt(26), bold=True, color=NAVY)
textbox(Inches(0.55), Inches(1.42), Inches(11.81), Inches(0.42),
        'Every feature below is built and verified, not planned. Right: an illustrative traditional-agency '
        'cost for the same scope \u2014 not a real quote, for contrast only.',
        Pt(12.5), color=GRAY)
textbox(Inches(0.55), Inches(7.14), Inches(8.00), Inches(0.30), 'BALKANEA MOBILE  \u2014  SANDBOX / TEST READINESS', Pt(9), color=FOOTER_GRAY)
textbox(Inches(11.78), Inches(7.14), Inches(1.00), Inches(0.30), '10', Pt(9), color=FOOTER_GRAY, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# Left column -- scope built
# ---------------------------------------------------------------------------
LEFT_X, LEFT_W = Inches(0.55), Inches(5.85)
TOP0 = Inches(1.98)

textbox(LEFT_X, TOP0, LEFT_W, Inches(0.22), 'SCOPE BUILT', Pt(10.5), bold=True, color=BROWN)

features = [
    ('Cross-platform app shell (iOS / Android / Web) + AI concierge (Nea)', False),
    ('850K-hotel content DB, region-filtered search', False),
    ('Real RateHawk booking flow \u2014 search / prebook / finish / status, all 7 certification scenarios, multi-room (real per-room guest configurator, not the flat-integer gap flagged earlier)', False),
    ('Real Bankart/NLB payment integration + voucher/invoice PDF retrieval', False),
    ('Auth: Google OAuth, native Apple Sign-In, guest checkout', False),
    ('Salesforce CRM sync, admin portal', False),
    ('Automatic booking emails \u2014 guest confirmation + business notification, each with a RateHawk PDF attached (voucher / invoice), shipped today via Resend', False),
    ('Push notifications', False),
    ('Voice AI (LiveKit/Retell) \u2014 built and working, shelved for a leaner V1', False),
]
multiline_box(LEFT_X, Emu(int(TOP0 + Inches(0.30))), LEFT_W, Inches(2.55),
              [(f'\u2022  {t}', False) for t, _ in features], Pt(9.3))

# Explicit callouts -- languages/currencies distinction, per Ray's ask
callout_top = Emu(int(TOP0 + Inches(2.95)))
rect = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, LEFT_X, callout_top, LEFT_W, Inches(0.95))
rect.adjustments[0] = 0.10
rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
rect.line.color.rgb = CARD_BORDER; rect.line.width = Pt(0.75)
rect.shadow.inherit = False
multiline_box(Emu(int(LEFT_X + Pt(10))), Emu(int(callout_top + Pt(8))), Emu(int(LEFT_W - Pt(20))), Inches(0.82), [
    ('Languages: 2 \u2014 English + Macedonian (Cyrillic), the actual translated UI.', True),
    ('Currencies: 9, across 11 countries \u2014 a separate locale/pricing display layer, not UI language. (11 countries is the destination-selector count, not a language count.)', False),
], Pt(9.3), line_gap=Pt(4))

# Codebase stats
stats_top = Emu(int(callout_top + Inches(1.10)))
textbox(LEFT_X, stats_top, LEFT_W, Inches(0.20), 'ACTUAL CODEBASE, ~10 WEEKS ELAPSED (first commit 6/18 \u2192 8/25)', Pt(9), bold=True, color=BROWN)
multiline_box(LEFT_X, Emu(int(stats_top + Inches(0.26))), LEFT_W, Inches(0.85), [
    ('Mobile: 17 screens \u00b7 12 components \u00b7 27 lib modules \u00b7 ~14,700 lines TS/TSX \u00b7 7 migrations \u00b7 86 commits', False),
    ('Backend: 12 endpoints \u00b7 12 lib modules \u00b7 ~4,200 lines JS \u00b7 73 commits', False),
], Pt(9.3), line_gap=Pt(4))

# ---------------------------------------------------------------------------
# Right column -- illustrative agency cost table
# ---------------------------------------------------------------------------
RIGHT_X, RIGHT_W = Inches(6.65), Inches(6.13)
textbox(RIGHT_X, TOP0, RIGHT_W, Inches(0.22), 'IF BUILT BY A TRADITIONAL AGENCY (ILLUSTRATIVE, NOT A REAL QUOTE)', Pt(10.5), bold=True, color=BROWN)

rows = [
    ('Discovery & UX', '$15\u201325K'),
    ('App shell / design system', '$10\u201315K'),
    ('Auth', '$8\u201312K'),
    ('Hotel search / content DB', '$20\u201330K'),
    ('AI concierge chat', '$25\u201340K'),
    ('Real booking flow (RateHawk)', '$35\u201355K'),
    ('Payment integration (Bankart)', '$25\u201340K'),
    ('Admin portal', '$8\u201312K'),
    ('Salesforce CRM', '$8\u201312K'),
    ('i18n', '$8\u201312K'),
    ('Voice AI (shelved)', '$15\u201325K'),
    ('QA across iOS/Android/Web', '$15\u201325K'),
]

table_top = Emu(int(TOP0 + Inches(0.30)))
n_rows = len(rows) + 1
row_h = Inches(0.205)
table_h = Emu(int(row_h * n_rows))
gtable = sp.add_table(n_rows, 2, RIGHT_X, table_top, RIGHT_W, table_h)
table = gtable.table
table.columns[0].width = Emu(int(RIGHT_W * 0.68))
table.columns[1].width = Emu(int(RIGHT_W * 0.32))

# Header row
hdr = [('Component', PP_ALIGN.LEFT), ('Blended $100\u2013180/hr', PP_ALIGN.RIGHT)]
for c, (text, align) in enumerate(hdr):
    cell = table.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    cell.margin_left = Pt(6); cell.margin_right = Pt(6); cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    tf = cell.text_frame
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(8.5); run.font.bold = True; run.font.color.rgb = WHITE
    run.font.name = 'Segoe UI Semibold'

for r, (label, cost) in enumerate(rows, start=1):
    for c, (text, align) in enumerate([(label, PP_ALIGN.LEFT), (cost, PP_ALIGN.RIGHT)]):
        cell = table.cell(r, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
        cell.margin_left = Pt(6); cell.margin_right = Pt(6); cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
        tf = cell.text_frame
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(8.7); run.font.color.rgb = NAVY if c == 0 else GRAY
        run.font.name = 'Segoe UI'

# Kill default pptx table style banding so our explicit fills show cleanly
tbl_el = table._tbl
tblPr = tbl_el.find(qn('a:tblPr'))
if tblPr is not None:
    tblPr.set('firstRow', '0')
    tblPr.set('bandRow', '0')

# Summary banner below the table
summary_top = Emu(int(table_top + table_h + Inches(0.14)))
rect2 = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, RIGHT_X, summary_top, RIGHT_W, Inches(1.30))
rect2.adjustments[0] = 0.08
rect2.fill.solid(); rect2.fill.fore_color.rgb = AMBER_FILL
rect2.line.fill.background()
rect2.shadow.inherit = False
multiline_box(Emu(int(RIGHT_X + Pt(10))), Emu(int(summary_top + Pt(8))), Emu(int(RIGHT_W - Pt(20))), Inches(1.15), [
    ('Subtotal ~$192K\u2013303K, +15\u201320% PM/coordination \u2192 ~$220K\u2013360K all-in.', True),
    ('Typically 4\u20137+ months with a small team. Nearshore/offshore: ~40\u201350% of that (~$110K\u2013180K), usually a longer timeline.', False),
], Pt(9.5), line_gap=Pt(5))

# Contrast banner
contrast_top = Emu(int(summary_top + Inches(1.44)))
rect3 = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, RIGHT_X, contrast_top, RIGHT_W, Inches(0.62))
rect3.adjustments[0] = 0.15
rect3.fill.solid(); rect3.fill.fore_color.rgb = GREEN_FILL
rect3.line.fill.background()
rect3.shadow.inherit = False
textbox(Emu(int(RIGHT_X + Pt(10))), Emu(int(contrast_top + Pt(10))), Emu(int(RIGHT_W - Pt(20))), Inches(0.42),
        'This build: ~10 weeks, ~$1,000 estimated engineering cost (see Cost slide).',
        Pt(11), bold=True, color=GREEN_TEXT)

prs.save(PATH)

# ---------------------------------------------------------------------------
# Move the new slide (currently last) into position INSERT_AT
# ---------------------------------------------------------------------------
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
new_el = slides[-1]
xml_slides.remove(new_el)
xml_slides.insert(INSERT_AT, new_el)
prs.save(PATH)
print('Inserted "What\u2019s Been Built" slide at index', INSERT_AT, '- total slides now:', len(prs.slides))
