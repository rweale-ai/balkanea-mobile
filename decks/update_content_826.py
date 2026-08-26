# -*- coding: utf-8 -*-
"""
2026-08-26 update: corrects three stale/wrong Risk items (multi-room, real-hotel
payments, RateHawk wiring -- all actually resolved/built as of 8/25), updates
Slides 1/3/4/9/10/12 text to match, and replaces the Risks slide's 4 cards with
the 4 genuinely-current risks. Reuses the card/badge/textbox pattern from
add_risks_slide.py so the visual style matches exactly.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)

# ---------------------------------------------------------------------------
# Part 1: targeted text replacements (preserve existing formatting/position)
# ---------------------------------------------------------------------------
REPLACEMENTS = {
    0: [  # Slide 1 - title
        (
            "RateHawk sandbox booking and payments are treated as working for this update -- see Risks for what that assumes. Cost model updated, certification and store submission plan attached.",
            "RateHawk sandbox booking and payments are now verified live end-to-end -- three real bookings processed today, including a fixed root-cause payment bug. See Risks for what's still genuinely open. Cost model updated, certification and store submission plan attached.",
        ),
    ],
    2: [  # Slide 3 - Search & Booking
        (
            "Search \u2192 prebook \u2192 create order \u2192 finish now runs successfully against RateHawk's certification test hotels (Conrad LA, Rosa Bell Motel) -- all 7 mandatory scenarios passed, real order IDs captured. Live-app wiring and the room-by-room guest picker UI are next -- see Risks and Plan.",
            "Search \u2192 prebook \u2192 create order \u2192 finish now runs successfully against RateHawk's certification test hotels (Conrad LA, Rosa Bell Motel) -- all 7 mandatory scenarios passed, real order IDs captured. Live-app wiring and the room-by-room guest picker UI are both built and verified -- three real bookings processed through the actual app today.",
        ),
    ],
    3: [  # Slide 4 - Payments
        (
            "Confirmed 8/20: a real test-card payment cleared through NLB via Bankart, from the app's booking flow to a synced Salesforce lead. Settles in MKD only -- see Risks for the one open item.",
            "Confirmed 8/20 (demo hotel) and 8/25 (real RateHawk hotel, after fixing a merchantTransactionId length bug): real test-card payments clear through NLB via Bankart, from the app's booking flow to a synced Salesforce lead. Settles in MKD -- see Risks for what's still open.",
        ),
        (
            "Hristijan's \u2018Balkanea Payment Bridge\u2019 plugin (signed HMAC links, Bankart Gateway V3 hosted card form, async Notify webhook) is live. Confirmed 8/20: a real test-card payment cleared through NLB, initiated from the actual mobile app booking flow -- 3DS, async postback, booking auto-confirm, and Salesforce sync all verified working end-to-end, settling in MKD (Bankart's only supported currency).",
            "Hristijan's \u2018Balkanea Payment Bridge\u2019 plugin (signed HMAC links, Bankart Gateway V3 hosted card form, async Notify webhook) is live. Confirmed 8/20 and re-confirmed 8/25 against a real RateHawk hotel, settling in MKD (Bankart's only supported currency, confirmed directly with Hristijan). 8/25 also found and fixed the actual root cause of intermittent failures: Hristijan's plugin appends its own timestamp suffix to our payment reference before submitting to Bankart, silently exceeding Bankart's 50-character limit on real-hotel bookings. Fixed app-side; charges now clear reliably.",
        ),
        (
            "Bankart's hosted card page isn't mobile-responsive today \u2014 input fields misalign on phone and tablet, iOS and Android. That's the one open item blocking a clean guest-facing checkout; on Hristijan's side to fix.",
            "Bankart's hosted card page still isn't mobile-responsive, and also doesn't collect billing address or a separate \u2018name on card\u2019 field (needed since the payer isn't always the guest). A full written spec -- Balkanea's brand colors, autofill support, and these missing fields -- was sent to Hristijan 8/25. Still his to build.",
        ),
    ],
    8: [  # Slide 9 - Architecture
        (
            "Bookings, profiles, payment_reference / payment_state (MKD). Migration live; real booking writes still pending the app-side RateHawk wiring (Plan, Phase 2).",
            "Bookings, profiles, payment_reference / payment_state (MKD), ratehawk_order_id. Live -- three real bookings written today via the fully-wired RateHawk flow. New: voucher/invoice PDF fetch and automated confirmation emails (guest + business team) via Resend.",
        ),
    ],
    9: [  # Slide 10 - Cost
        (
            "Estimated $400 AI Cost (excludes the unscoped multi-room build -- see Risks)",
            "Estimated $1,000 engineering cost (Claude Code, if metered -- near-zero if on subscription). Reflects 8/25's additional real-money payment debugging, reliability fixes, and the new voucher/invoice/email feature; multi-room is built, not excluded.",
        ),
    ],
    12: [  # Slide 13 - Still Open
        (
            "From Ray: Apple Services ID/JWT, Google Play account creation. From Hristijan: Bankart mobile-responsive fix. From Claude: RateHawk app-wiring and the multi-room build per the Plan slide.",
            "From Ray: Apple Services ID/JWT, Google Play account creation. From Hristijan: Bankart mobile-responsive fix + billing-address/cardholder-name fields (spec sent 8/25). From Claude: RateHawk app-wiring and multi-room build are both done and verified -- next is RateHawk's certification questionnaire.",
        ),
        (
            "Bankart's hosted card page isn't mobile-responsive today -- input fields misalign on phone and tablet, iOS and Android. On Hristijan's side; the one item blocking a clean guest-facing checkout.",
            "Bankart's hosted card page isn't mobile-responsive today, and is missing billing-address and cardholder-name fields. Written spec sent to Hristijan 8/25 (design tokens, autofill attributes, field list). On his side to build.",
        ),
    ],
}

changed = 0
for idx, pairs in REPLACEMENTS.items():
    slide = prs.slides[idx]
    for old, new in pairs:
        found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full = ''.join(r.text for r in para.runs)
                if old in full or old.strip() == full.strip():
                    # Replace by rewriting the first run's text and clearing the rest,
                    # preserving the first run's formatting.
                    if para.runs:
                        para.runs[0].text = new
                        for r in para.runs[1:]:
                            r.text = ''
                        found = True
                        changed += 1
                        break
            if found:
                break
        if not found:
            print(f"WARNING: text not found on slide index {idx}: {old[:80]}...")

print(f"Replaced {changed} text blocks across slides 1/3/4/9/10/12")

# ---------------------------------------------------------------------------
# Part 2: rebuild the Risks slide (index 11) with the 4 current risks
# ---------------------------------------------------------------------------
RISKS_IDX = 11
slide = prs.slides[RISKS_IDX]

# Remove every shape on the slide (background stays -- it's a slide property,
# not a shape).
for shape in list(slide.shapes):
    shape._element.getparent().remove(shape._element)

sp = slide.shapes

CREAM = RGBColor(0xFA, 0xF7, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE7, 0xE0, 0xD3)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x65, 0x58)
FOOTER_GRAY = RGBColor(0x9B, 0x93, 0x84)
BROWN = RGBColor(0x7A, 0x3E, 0x0F)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD9)
AMBER_TEXT = RGBColor(0xB8, 0x79, 0x0A)
GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_TEXT = RGBColor(0x2E, 0x6B, 0x57)


def textbox(left, top, width, height, text, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = sp.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.name = 'Segoe UI Semibold' if bold else 'Segoe UI'
    run.font.color.rgb = color
    return box


def badge(rect_ref, label, positive=True):
    w = Emu(int(Emu(Pt(5.3)) * len(label) + Emu(Pt(16))))
    h = Inches(0.20)
    left = Emu(int(Emu(rect_ref.left) + Emu(rect_ref.width) - w - Pt(8)))
    top = Emu(int(Emu(rect_ref.top) + Pt(7)))
    b = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    b.adjustments[0] = 0.5
    b.fill.solid()
    b.fill.fore_color.rgb = GREEN_FILL if positive else AMBER_FILL
    b.line.fill.background()
    b.shadow.inherit = False
    tf = b.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(7.5)
    run.font.bold = True
    run.font.name = 'Segoe UI Semibold'
    run.font.color.rgb = GREEN_TEXT if positive else AMBER_TEXT
    return b


def card(left, top, width, height, title, desc, badge_label, badge_positive):
    rect = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.09
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = CARD_BORDER
    rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    pad = Pt(10)
    title_top = Emu(int(top + Inches(0.12)))
    textbox(Emu(int(left + pad)), title_top, Emu(int(width - 2 * pad - Inches(1.6))), Inches(0.22), title, Pt(11.5), bold=True, color=NAVY)
    desc_top = Emu(int(title_top + Inches(0.28)))
    textbox(Emu(int(left + pad)), desc_top, Emu(int(width - 2 * pad)), Emu(int(top + height - desc_top - Pt(6))), desc, Pt(9.5), color=GRAY)
    if badge_label:
        badge(rect, badge_label, positive=badge_positive)
    return rect


# Header
textbox(Inches(0.55), Inches(0.50), Inches(12.23), Inches(0.35), '11  \u2014  RISKS', Pt(12), bold=True, color=BROWN)
textbox(Inches(0.55), Inches(0.82), Inches(12.23), Inches(0.70),
        "What's actually still open, now that payments and the RateHawk flow are verified live", Pt(27), bold=True, color=NAVY)
textbox(Inches(0.55), Inches(1.55), Inches(11.81), Inches(0.55),
        'Multi-room booking, real-hotel payments, and RateHawk\u2019s live wiring -- all previously flagged here -- are '
        'now built, fixed, and verified as of 8/25. These are the risks that remain.',
        Pt(13.5), color=GRAY)
textbox(Inches(0.55), Inches(7.14), Inches(8.00), Inches(0.30), 'BALKANEA MOBILE  \u2014  SANDBOX / TEST READINESS', Pt(9), color=FOOTER_GRAY)
textbox(Inches(11.78), Inches(7.14), Inches(1.00), Inches(0.30), '12', Pt(9), color=FOOTER_GRAY, align=PP_ALIGN.RIGHT)

CARD_X, CARD_W = Inches(0.55), Inches(12.23)
TOP0 = Inches(2.35)
CARD_H = Inches(1.02)
GAP = Inches(0.15)

card(CARD_X, TOP0, CARD_W, CARD_H,
    "RateHawk's sandbox documents (voucher + invoice) return a blank placeholder, not real content",
    "Fetching both via RateHawk's real document endpoints returns byte-for-byte identical PDFs -- embedded "
    "metadata reads 'Acrobat Distiller 6.0 (Windows)', dated 2006-03-06, clearly a generic template, not a "
    "generated document. Not a bug in our integration (request/response plumbing verified correct) -- needs "
    "re-verification once production RateHawk credentials are active; sandbox may simply have no real content to render.",
    'SANDBOX LIMITATION', False)

card(CARD_X, Emu(int(TOP0 + CARD_H + GAP)), CARD_W, CARD_H,
    'A reliability bug caused three duplicate real charges today, before being fixed',
    "The booking screen could get stuck indefinitely on \u2018Confirming with the hotel\u2019 if a step failed "
    "after payment had already captured. Because real RateHawk confirmation genuinely takes 90-140+ seconds, "
    "this read as frozen; force-quitting and retrying created a brand-new charge each time. Result: three real "
    "Bankart charges for one intended test booking -- harmless only because it was a test card. Now fixed "
    "(proper error states + a real progress bar); the lesson is this exact bug class is what causes real "
    "duplicate charges against a real guest.",
    'FIXED TODAY', True)

card(CARD_X, Emu(int(TOP0 + 2 * (CARD_H + GAP))), CARD_W, CARD_H,
    "Bankart's hosted card page still isn't mobile-responsive",
    "Input fields still misalign on phone and tablet, and the page has no billing-address or independent "
    "'name on card' field (a guest can pay for someone else's booking). A full written spec -- brand colors, "
    "autofill attributes, missing fields -- was sent to Hristijan 8/25. Not yet built on his side.",
    'OPEN \u2014 HRISTIJAN', False)

card(CARD_X, Emu(int(TOP0 + 3 * (CARD_H + GAP))), CARD_W, CARD_H,
    "Business-team booking notification email is hardcoded to a placeholder address",
    "New today: confirmed real RateHawk bookings now trigger automated emails -- the guest gets their voucher "
    "PDF, the business team gets the invoice PDF (Resend). The business recipient is currently hardcoded to "
    "Ray's own address as a placeholder, not a real shared team inbox -- needs to move before this is "
    "production-ready.",
    'PLACEHOLDER', False)

prs.save(PATH)
print("Saved Risks slide rebuild and all text replacements.")
