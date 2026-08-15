import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PATH = 'Balkanea-Mobile-Sandbox-Readiness-Update.pptx'
prs = Presentation(PATH)

# Drop a previously-built architecture slide if this script is being re-run (always the last slide).
if len(prs.slides) == 9:
    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(list(xml_slides)[8])

CREAM = RGBColor(0xFA, 0xF7, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE7, 0xE0, 0xD3)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x65, 0x58)
FOOTER_GRAY = RGBColor(0x9B, 0x93, 0x84)
BROWN = RGBColor(0x7A, 0x3E, 0x0F)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD9)
AMBER_TEXT = RGBColor(0xB8, 0x79, 0x0A)
GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_TEXT = RGBColor(0x2E, 0x6B, 0x57)
LINE_GRAY = RGBColor(0xAF, 0xA6, 0x94)

layout = prs.slides[4].slide_layout
new_slide = prs.slides.add_slide(layout)
new_slide.background.fill.solid()
new_slide.background.fill.fore_color.rgb = CREAM
sp = new_slide.shapes


def textbox(left, top, width, height, text, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = sp.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.name = 'Segoe UI Semibold' if bold else 'Segoe UI'
    run.font.color.rgb = color
    return box


def card(left, top, width, height, dashed=False):
    rect = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.09
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = CARD_BORDER
    rect.line.width = Pt(0.75)
    if dashed:
        ln = rect.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
        rect.line.color.rgb = AMBER_TEXT
        rect.line.width = Pt(1.25)
    rect.shadow.inherit = False
    return rect


def badge(rect_ref, label, positive=True):
    w = Emu(int(Emu(Pt(5.3)) * len(label) + Emu(Pt(16))))
    h = Inches(0.20)
    left = Emu(int(Emu(rect_ref.left) + Emu(rect_ref.width) - w - Pt(8)))
    top = Emu(int(Emu(rect_ref.top) + Pt(7)))
    b = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    b.adjustments[0] = 0.5
    b.fill.solid()
    b.fill.fore_color.rgb = GREEN_FILL if positive else AMBER_FILL
    b.line.fill.background()
    b.shadow.inherit = False
    tf = b.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(7.5)
    run.font.bold = True
    run.font.name = 'Segoe UI Semibold'
    run.font.color.rgb = GREEN_TEXT if positive else AMBER_TEXT
    return b


def component(left, top, width, height, title, desc, badge_label=None, badge_positive=True, dashed=False):
    rect = card(left, top, width, height, dashed=dashed)
    pad = Pt(10)
    # Badged cards need the title pushed below the badge chip to avoid overlap.
    title_top = Emu(int(top + (Inches(0.32) if badge_label else Pt(9))))
    textbox(Emu(int(left + pad)), title_top, Emu(int(width - 2 * pad)), Inches(0.22), title, Pt(11.5), bold=True,
            color=(AMBER_TEXT if dashed else NAVY))
    desc_top = Emu(int(title_top + Inches(0.26)))
    textbox(Emu(int(left + pad)), desc_top, Emu(int(width - 2 * pad)), Emu(int(top + height - desc_top - Pt(6))), desc, Pt(9), color=GRAY)
    if badge_label:
        badge(rect, badge_label, positive=badge_positive)
    return rect


def right_mid(shape):
    return Emu(int(shape.left + shape.width)), Emu(int(shape.top + shape.height / 2))


def left_mid(shape):
    return Emu(shape.left), Emu(int(shape.top + shape.height / 2))


def bottom_mid(shape):
    return Emu(int(shape.left + shape.width / 2)), Emu(int(shape.top + shape.height))


def top_mid(shape):
    return Emu(int(shape.left + shape.width / 2)), Emu(shape.top)


def arrow(x1, y1, x2, y2, dashed=False):
    conn = sp.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = AMBER_TEXT if dashed else LINE_GRAY
    conn.line.width = Pt(1.5 if dashed else 1.1)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if dashed:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return conn


def flow_label(cx, cy, text, color=FOOTER_GRAY):
    w = Inches(1.3)
    box = textbox(Emu(int(cx - w / 2)), Emu(int(cy - Inches(0.09))), w, Inches(0.18), text, Pt(8), color=color, align=PP_ALIGN.CENTER)
    box.fill.solid()
    box.fill.fore_color.rgb = CREAM
    box.line.fill.background()
    return box


# ---------------------------------------------------------------------------
# Header
# ---------------------------------------------------------------------------
textbox(Inches(0.55), Inches(0.50), Inches(12.23), Inches(0.35), '05  —  ARCHITECTURE', Pt(12), bold=True, color=BROWN)
textbox(Inches(0.55), Inches(0.82), Inches(12.23), Inches(0.55), 'Planned production architecture', Pt(27), bold=True, color=NAVY)
textbox(Inches(0.55), Inches(1.40), Inches(11.81), Inches(0.40),
        'One box is still open — payments. Everything else here is either live today or a scoped, unblocked build.',
        Pt(13.5), color=GRAY)
textbox(Inches(0.55), Inches(7.14), Inches(8.00), Inches(0.30), 'BALKANEA MOBILE  —  SANDBOX / TEST READINESS', Pt(9), color=FOOTER_GRAY)
page_num_box = textbox(Inches(11.78), Inches(7.14), Inches(1.00), Inches(0.30), '?', Pt(9), color=FOOTER_GRAY, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# Columns
# ---------------------------------------------------------------------------
COL1_X, COL1_W = Inches(0.55), Inches(2.55)
COL2_X, COL2_W = Inches(3.44), Inches(2.55)
COL3_X, COL3_W = Inches(6.33), Inches(2.55)
COL4_X, COL4_W = Inches(9.22), Inches(3.56)

HEADER_TOP = Inches(1.95)
for x, w, label in [(COL1_X, COL1_W, 'CLIENT'), (COL2_X, COL2_W, 'BACKEND'), (COL3_X, COL3_W, 'DATA'), (COL4_X, COL4_W, 'PARTNERS')]:
    textbox(x, HEADER_TOP, w, Inches(0.2), label, Pt(9), bold=True, color=FOOTER_GRAY, align=PP_ALIGN.CENTER)

BODY_TOP = Inches(2.24)

# --- Client ---
mobile_h = Inches(1.35)
mobile = component(COL1_X, BODY_TOP, COL1_W, mobile_h,
    'Balkanea Mobile App',
    'Expo / React Native — iOS, Android, Web. Nea AI (Claude Haiku) search advisor. Booking flow: search → lock → pay → confirm.')

auth_top = Emu(int(BODY_TOP + mobile_h + Inches(0.24)))
auth_h = Inches(1.05)
auth = component(COL1_X, auth_top, COL1_W, auth_h,
    'Supabase Auth',
    'Google (live) · Apple (native code shipped, provider secret pending) · Phone OTP · Guest.')

# --- Backend (hub) ---
backend_h = Inches(1.95)
backend_top = Emu(int(BODY_TOP + (mobile_h + Inches(0.24) + auth_h) / 2 - backend_h / 2))
backend = component(COL2_X, backend_top, COL2_W, backend_h,
    'balkanea-lead-webhook',
    'Vercel serverless (Node.js). Hotel search proxy, RateHawk mobile workflow (SERP / Prebook / create-order), payment gateway, Salesforce sync.')

# --- Data ---
hotels_h = Inches(1.35)
hotels_db = component(COL3_X, BODY_TOP, COL3_W, hotels_h,
    'balkanea_hotels_sbx',
    '850,218 hotels / 3.9M rooms, all 10 countries. Region-filtered search, daily refresh from RateHawk dump.',
    badge_label='LIVE', badge_positive=True)

bookings_top = Emu(int(BODY_TOP + hotels_h + Inches(0.24)))
bookings_h = Inches(1.35)
bookings_db = component(COL3_X, bookings_top, COL3_W, bookings_h,
    'balkanea_sbx',
    'Bookings, profiles, payment_reference / payment_state. Migration ready, on hold pending the payment decision.')

# --- Partners ---
ratehawk_h = Inches(1.35)
ratehawk = component(COL4_X, BODY_TOP, COL4_W, ratehawk_h,
    'RateHawk',
    'Separate mobile Office (contract "balkanea-mobile", confirmed 8/14). Sandbox key incoming → certification → production key.',
    badge_label='IN PROGRESS', badge_positive=True)

nlb_top = Emu(int(BODY_TOP + ratehawk_h + Inches(0.22)))
nlb_h = Inches(1.05)
nlb = component(COL4_X, nlb_top, COL4_W, nlb_h,
    'Payment gateway (NLB-direct)',
    'No WooCommerce, per Ray. Real spec not yet in hand — asked Hristijan/Ivan. Nothing here is built.',
    badge_label='TBD', badge_positive=False, dashed=True)

sf_top = Emu(int(nlb_top + nlb_h + Inches(0.22)))
sf_h = Inches(0.95)
salesforce = component(COL4_X, sf_top, COL4_W, sf_h,
    'Salesforce CRM',
    'Leads, booking sync, agent escalation.')

# ---------------------------------------------------------------------------
# Connectors — straight fan-out spokes from a single hub point, not auto-routed elbows
# ---------------------------------------------------------------------------
# Mobile -> Auth (client SDK, direct — bypasses the backend)
x1, y1 = bottom_mid(mobile)
x2, y2 = top_mid(auth)
arrow(x1, y1, x2, y2)
flow_label(x1, Emu(int((y1 + y2) / 2)), 'client SDK')

# Mobile -> Backend
x1, y1 = right_mid(mobile)
x2, y2 = left_mid(backend)
arrow(x1, y1, x2, y2)

# Backend hub point (right edge, vertical center) fans out to every Data/Partner box
hub_x, hub_y = right_mid(backend)
for target in [hotels_db, bookings_db, ratehawk, salesforce]:
    tx, ty = left_mid(target)
    arrow(hub_x, hub_y, tx, ty)
tx, ty = left_mid(nlb)
arrow(hub_x, hub_y, tx, ty, dashed=True)

# ---------------------------------------------------------------------------
# Legend
# ---------------------------------------------------------------------------
legend_top = Inches(6.68)
leg_line = sp.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.55), Emu(int(legend_top + Inches(0.09))), Inches(0.95), Emu(int(legend_top + Inches(0.09))))
leg_line.line.color.rgb = AMBER_TEXT
leg_line.line.width = Pt(1.5)
ln = leg_line.line._get_or_add_ln()
ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
textbox(Inches(1.05), legend_top, Inches(7.0), Inches(0.2), '= not yet built, pending the NLB-direct spec from Hristijan/Ivan', Pt(9), color=FOOTER_GRAY)

prs.save(PATH)
print('Slide added. Total slides now:', len(prs.slides))
