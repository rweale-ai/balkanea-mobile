# -*- coding: utf-8 -*-
"""
Content-only edits (no slide insert/delete) to Balkanea-Mobile-Sandbox-Readiness-Update.pptx,
reflecting status as of 2026-08-25:
 - RateHawk sandbox key 973 (replacing blocked 954), verified; real booking mechanics
   proven end-to-end incl. 7/7 certification scenarios (multiroom included).
 - Payments: MKD-only settlement confirmed; treated as working per Ray's direction for
   this update (open settlement issue moved to the new Risks slide, not stated here).
 - Multi-room booking: treated as working per Ray's direction on this slide (mechanics
   proven); the app-side UI/data-model gap is NOT claimed resolved here, and is the lead
   item on the new Risks slide.
 - Proxy/AWS slides re-scoped from "sandbox unblock" to "production certification prerequisite"
   (sandbox needs no IP whitelisting -- confirmed 8/24).
 - Planning slide rewritten into 4 phases explicitly covering RateHawk certification,
   testing, and both store submissions.
 - Architecture slide's stale boxes (sandbox key incoming, payment bridge status, stage-
   server proxy) brought current.
 - Cost slide 10 gets a short AWS-relay cost note; slide 11 numbers refreshed to match
   the recalculated cost model (Excel-confirmed values).
 - "Still Open" slide re-scoped to genuinely-still-open operational loose ends (Apple/
   Google Play accounts, Bankart mobile-responsive page, admin portal parity gap) since
   the big strategic items moved to the new Risks slide.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)

GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_FONT = RGBColor(0x2E, 0x6B, 0x57)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD9)
AMBER_FONT = RGBColor(0xB8, 0x79, 0x0A)


def set_text(shape, new_text):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    proto = p.runs[0]
    proto_size = proto.font.size
    proto_bold = proto.font.bold
    proto_name = proto.font.name
    proto_color = proto.font.color.rgb
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    for extra_p in list(tf.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)
    run = p.add_run()
    run.text = new_text
    run.font.size = proto_size
    run.font.bold = proto_bold
    run.font.name = proto_name
    run.font.color.rgb = proto_color
    return run


def recolor_badge(fill_shape, label_shape, text, positive):
    fill_shape.fill.solid()
    fill_shape.fill.fore_color.rgb = GREEN_FILL if positive else AMBER_FILL
    set_text(label_shape, text)
    label_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = GREEN_FONT if positive else AMBER_FONT


# =====================================================================
# SLIDE 0 -- TITLE
# =====================================================================
s = prs.slides[0]
sh = s.shapes
set_text(sh[1], "Cost, risk, and the path to submission")
set_text(
    sh[2],
    "RateHawk sandbox booking and payments are treated as working for this update "
    "-- see Risks for what that assumes. Cost model updated, certification and store "
    "submission plan attached.",
)
set_text(
    sh[3],
    "Stage:  TestFlight (v1.0.7, build 11)      Systems:  RateHawk (sandbox, cert pending) "
    "\u00b7 Bank checkout (MKD, live) \u00b7 Supabase auth      Prepared:  August 2026 (update)",
)

# =====================================================================
# SLIDE 2 -- SEARCH & BOOKING
# =====================================================================
s = prs.slides[2]
sh = s.shapes

set_text(sh[1], "RateHawk sandbox is live; full booking flow proven end-to-end")
set_text(
    sh[2],
    "Sandbox key 973 replaced the blocked key 954 and returns real search results across "
    "all sandbox regions. The complete booking sequence -- search, prebook, create order, "
    "finish -- has been proven against RateHawk's own certification test hotels, including "
    "every mandatory certification scenario.",
)

# Card 1: BLOCKED -> RESOLVED (key 973)
recolor_badge(sh[6], sh[7], "RESOLVED", True)
set_text(sh[8], "Sandbox key 973 -- issued and verified 8/24, replacing blocked key 954")
set_text(
    sh[9],
    "RateHawk confirmed the 401s on key 954 were a technical issue on their side, not IP "
    "whitelisting -- a new key (973) was issued and verified with real search results. "
    "Sandbox itself needs no IP whitelisting at all; that only applies at production "
    "certification (see Risks).",
)

# Card 4: Open -> CONFIRMED (booking mechanics proven; app wiring + UI still ahead)
recolor_badge(sh[16], sh[17], "CONFIRMED", True)
set_text(sh[18], "Full booking flow proven end-to-end in sandbox, including multi-room")
set_text(
    sh[19],
    "Search \u2192 prebook \u2192 create order \u2192 finish now runs successfully against "
    "RateHawk's certification test hotels (Conrad LA, Rosa Bell Motel) -- all 7 mandatory "
    "certification scenarios passed, including a 2-room booking with mixed child ages, with "
    "real order IDs captured for the certification submission. Wiring this into the live "
    "app and building the room-by-room guest picker UI are the next build items -- see the "
    "Risks and Plan slides; treated as working for this update per Ray's direction.",
)

# =====================================================================
# SLIDE 3 -- PAYMENTS
# =====================================================================
s = prs.slides[3]
sh = s.shapes

set_text(sh[1], "Bank checkout works end-to-end -- confirmed live with NLB, in MKD")
set_text(
    sh[2],
    "Confirmed 8/20: a real test-card payment cleared through NLB via Bankart, from the "
    "app's booking flow to a synced Salesforce lead. Bankart's merchant account is MKD-only "
    "(confirmed by Hristijan) -- the app settles in MKD regardless of the guest's selected "
    "display currency. Treated as working for this update; see Risks for the one open item.",
)

set_text(sh[13], "Balkanea Payment Bridge -- built and confirmed working end-to-end, MKD settlement")
set_text(
    sh[14],
    "Hristijan's \u2018Balkanea Payment Bridge\u2019 plugin (signed HMAC links, Bankart "
    "Gateway V3 hosted card form, async Notify webhook) is live. Confirmed 8/20: a real "
    "test-card payment cleared through NLB, initiated from the actual mobile app booking "
    "flow -- 3DS, async postback, booking auto-confirm, and Salesforce sync all verified "
    "working end-to-end, settling in MKD (Bankart's only supported currency).",
)

# =====================================================================
# SLIDE 4 -- PROXY: Proposal vs AWS  (re-scope: production cert prerequisite, not sandbox unblock)
# =====================================================================
s = prs.slides[4]
sh = s.shapes
set_text(
    sh[2],
    "Scoped to the static-IP RateHawk requires for PRODUCTION certification -- sandbox "
    "itself needs no whitelisting (confirmed 8/24). Both columns are real options.",
)

# =====================================================================
# SLIDE 5 -- ALIGNMENT & RECOMMENDATION
# =====================================================================
s = prs.slides[5]
sh = s.shapes
set_text(
    sh[11],
    "For the actual requirement -- a static IP RateHawk's production certification will "
    "accept -- the AWS build is live, verified, and ~$5\u20138/mo today, with no further "
    "engineering or coordination needed. Ivan's fuller proposal is well-designed for a "
    "mature, high-traffic system, but adopting it as specified means more infrastructure, "
    "more cost, and an unresolved question of who builds and runs it -- none of which the "
    "certification requirement itself needs. Fastest, lowest-risk path to a certification "
    "submission: use what's already built and verified. Revisit Redis/queueing/the full "
    "gateway abstraction once real production traffic justifies it.",
)

# =====================================================================
# SLIDE 7 -- PLANNING  (full rewrite: certification + build + test + submit)
# =====================================================================
s = prs.slides[7]
sh = s.shapes

set_text(sh[1], "Certification, build, test, submit -- phased")
set_text(
    sh[2],
    "Claude does all technical implementation; testing and store-asset design draw on "
    "existing team time. RateHawk's certification review and the app store review clocks "
    "are the two timelines genuinely outside our control.",
)

set_text(sh[7], "PHASE 1 - CERTIFY")
set_text(sh[8], "Ours to finish, RateHawk's to review")
set_text(
    sh[9],
    "Answer RateHawk's Pre-Certification Checklist in full (payment type = deposit, RPM "
    "limits per endpoint, static-data sync method, price-parsing field, booking-success "
    "signal, full error-handling matrix), submit the AWS static IP for whitelisting, and "
    "hand back the real order IDs already captured from the 7 passed test scenarios "
    "(Conrad LA + Rosa Bell Motel). RateHawk's own review time after submission isn't ours "
    "to set.",
)

set_text(sh[12], "PHASE 2 - BUILD")
set_text(sh[13], "~2\u20133 weeks \u00b7 Claude-implemented")
set_text(
    sh[14],
    "Replace lib/ratehawk.ts's simulated lockRoom/reconfirmBooking with the real prebook "
    "\u2192 charge \u2192 finish sequence, reconciling the fact that RateHawk charges "
    "Balkanea's own deposit balance while Bankart charges the guest's card -- two systems "
    "that can now disagree. Build the room-by-room guest configurator that search, room-"
    "selection, and booking currently have no UI for at all (a genuinely multi-file "
    "feature, not a small addition). Apple's OAuth secret and the Google Play Developer "
    "account (both Ray-owned) proceed in parallel, not build-gated.",
)

set_text(sh[17], "PHASE 3 - TEST")
set_text(sh[18], "~1\u20132 weeks, parallel across existing testers")
set_text(
    sh[19],
    "Auth, search, and single-room booking + payment are already TestFlight-tested. New "
    "coverage once Phase 2 lands: multi-room booking end-to-end, RateHawk's full error-"
    "handling matrix (booking_form_expired, rate_not_found, soldout, book_limit, 3ds -- not "
    "just the happy path), and a repeat real-money payment run to confirm settlement is "
    "clean before certification submission. Store assets (screenshots, listing copy) in "
    "parallel -- ~2\u20133 days on top of already-built UI.",
)

set_text(sh[22], "PHASE 4 - SUBMIT")
set_text(sh[23], "~1 month+, three review clocks running in parallel")
set_text(
    sh[24],
    "Apple review ~3\u20135 days (real travel payments = higher scrutiny; budget one "
    "rejection-and-resubmit round). Google Play's 14-day mandatory closed-testing period "
    "is the long pole, and its clock can't start until the Play Developer account exists. "
    "RateHawk moves sandbox \u2192 production keys only after certification sign-off, on "
    "its own timeline, running alongside store review rather than before it.",
)

# =====================================================================
# SLIDE 8 -- ARCHITECTURE
# =====================================================================
s = prs.slides[8]
sh = s.shapes

set_text(
    sh[2],
    "Payments and RateHawk booking are treated as live for this update (see Risks for what "
    "that assumes); the admin portal's WooCommerce-parity gap is the one box still "
    "genuinely open below.",
)

set_text(
    sh[27],
    "Separate mobile Office (contract \u201cbalkanea-mobile\u201d). Sandbox key 973 live "
    "and verified -- booking mechanics proven end-to-end against certification test "
    "hotels. Certification checklist submission is next, then production key.",
)

set_text(
    sh[24],
    "Bookings, profiles, payment_reference / payment_state (MKD). Migration live; real "
    "booking writes still pending the app-side RateHawk wiring (Plan, Phase 2).",
)

set_text(
    sh[31],
    "Confirmed 8/20: real end-to-end card payment cleared through NLB via Bankart's hosted "
    "card form, from the app's own booking flow, settling in MKD. Refunds still unbuilt.",
)
recolor_badge(sh[29], sh[32], "LIVE", True)

set_text(sh[48], "AWS Relay (static IP)")
set_text(
    sh[49],
    "AWS box in eu-central-1 (Frankfurt) -- static Elastic IP (63.185.27.197), live and "
    "verified. Required for RateHawk's production certification; sandbox itself needs no "
    "whitelisting.",
)
recolor_badge(sh[47], sh[50], "LIVE", True)

# =====================================================================
# SLIDE 9 -- COST (implementation)
# =====================================================================
s = prs.slides[9]
sh = s.shapes

set_text(sh[9], "Estimated $400 AI Cost (excludes the unscoped multi-room build -- see Risks)")

set_text(
    sh[29],
    "Sandbox running now at ~$25/mo (Pro plan, Medium compute). Production baseline "
    "~$75/mo once live; HA read replica deferred until traffic justifies it. Backups "
    "skipped. Plus the AWS static-IP relay now required for RateHawk production "
    "certification: ~$8/mo base, ~$22/mo with HA -- both in the updated cost model.",
)

# =====================================================================
# SLIDE 10 -- COST (running cost vs margin) -- refreshed to match the recalculated model
# =====================================================================
s = prs.slides[10]
sh = s.shapes

set_text(sh[12], "$1,018")
set_text(sh[17], "~$6,759")
set_text(sh[18], "Running cost = ~13% of margin revenue")

set_text(sh[26], "$8,555")
set_text(sh[31], "~$69,205")
set_text(sh[32], "Running cost = ~11% of margin revenue")

set_text(
    sh[33],
    "Same story at both scales -- compute/infrastructure (now including the AWS static-IP "
    "relay required for RateHawk certification) is a small, roughly flat share of margin "
    "(~11\u201313%). Confirms Luke's own read from July 21: marketing/adoption is the real "
    "bottleneck, not running cost. On Haiku 4.5 instead of Sonnet, running cost drops to "
    "~$518/mo (10's) and ~$3,555/mo (100's) -- further shrinking that share, at a quality "
    "trade-off worth testing before committing.",
)

# =====================================================================
# SLIDE 11 -- STILL OPEN  (re-scoped to genuinely-open operational loose ends;
#                            the big strategic items moved to the new Risks slide)
# =====================================================================
s = prs.slides[11]
sh = s.shapes

set_text(sh[1], "What's still genuinely open, and who owns it")
set_text(sh[2], "The strategic risks are on the Risks slide. This is the operational punch list.")

set_text(sh[6], "Apple + Google Play accounts")
set_text(
    sh[7],
    "Apple: native Sign-in code shipped 8/5, still needs a Services ID + signed JWT from "
    "Ray's own Apple Developer account. Google Play: developer account still not created "
    "(confirmed not started as of the July 24 call) -- account verification itself takes "
    "calendar time independent of app readiness. Both Ray-owned, not build-gated.",
)

set_text(sh[9], "Bankart mobile-responsive checkout")
set_text(
    sh[10],
    "Bankart's hosted card page isn't mobile-responsive today -- input fields misalign on "
    "phone and tablet, iOS and Android. On Hristijan's side; the one item blocking a clean "
    "guest-facing checkout.",
)

set_text(sh[12], "Admin portal WooCommerce-parity gap")
set_text(
    sh[13],
    "admin-payments.html covers bookings + payment state, with Capture/Void stubbed -- but "
    "doesn't yet cover WooCommerce's refunds/cancellations/invoicing. Gap inventoried, not "
    "yet built.",
)

set_text(
    sh[15],
    "deliverables to unblock this week",
)
set_text(
    sh[16],
    "From Ray: Apple Services ID/JWT, Google Play account creation. From Hristijan: "
    "Bankart mobile-responsive fix. From Claude: RateHawk app-wiring and the multi-room "
    "build per the Plan slide.",
)

set_text(
    sh[17],
    "Store submission, once ready: no Google Play account yet \u00b7 Apple review ~3\u20135 "
    "days \u00b7 Google's 14-day closed-testing minimum \u00b7 RateHawk certification runs "
    "in parallel -- see Plan and Risks.",
)

prs.save(PATH)
print("Content edits saved. Slide count:", len(prs.slides))
