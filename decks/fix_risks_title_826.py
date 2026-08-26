# -*- coding: utf-8 -*-
"""Fixes the Risks slide (index 11) title overflowing to 2 lines and
overlapping the subtitle -- shortens it to fit one line like the rest of
the deck's slide titles."""
from pptx import Presentation

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)
slide = prs.slides[11]

OLD = "What's actually still open, now that payments and the RateHawk flow are verified live"
NEW = "What's still open, now that payments and RateHawk are live"

fixed = False
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        full = ''.join(r.text for r in para.runs)
        if OLD in full:
            para.runs[0].text = NEW
            for r in para.runs[1:]:
                r.text = ''
            fixed = True

prs.save(PATH)
print("Fixed:", fixed)
