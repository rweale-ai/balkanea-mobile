# -*- coding: utf-8 -*-
"""
Full clean renumber pass, fixing the pre-existing drift (section headers skipped
"04", footers were off by one from slide 5 onward) as well as accounting for the
newly-inserted Risks slide.

Convention: slide at 0-based index si (si >= 1; si=0 is the title slide, untouched)
gets header f'{si:02d}' and footer str(si+1) -- i.e. footer matches the slide's
actual 1-based position in the deck.
"""
import re
from pptx import Presentation
from pptx.util import Inches

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)

HEADER_RE = re.compile(r'^(\d{2})(\s{2}\u2014\s{2}.+)$')

for si in range(1, len(prs.slides)):
    s = prs.slides[si]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()

        m = HEADER_RE.match(t)
        if m and len(t) < 40:
            new_text = f'{si:02d}' + m.group(2)
            if new_text != t:
                p = sh.text_frame.paragraphs[0]
                proto = p.runs[0]
                for r in list(p.runs):
                    r._r.getparent().remove(r._r)
                run = p.add_run()
                run.text = new_text
                run.font.size = proto.font.size
                run.font.bold = proto.font.bold
                run.font.name = proto.font.name
                run.font.color.rgb = proto.font.color.rgb
                print(f'slide {si}: header -> {new_text!r}')
            continue

        if sh.left is not None and sh.top is not None:
            if sh.left > Inches(11) and sh.top > Inches(7) and t.isdigit():
                new_text = str(si + 1)
                if new_text != t:
                    p = sh.text_frame.paragraphs[0]
                    proto = p.runs[0]
                    for r in list(p.runs):
                        r._r.getparent().remove(r._r)
                    run = p.add_run()
                    run.text = new_text
                    run.font.size = proto.font.size
                    run.font.bold = proto.font.bold
                    run.font.name = proto.font.name
                    run.font.color.rgb = proto.font.color.rgb
                    print(f'slide {si}: footer page -> {new_text!r}')

prs.save(PATH)
print('Saved. Total slides:', len(prs.slides))
