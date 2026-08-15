from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PATH = 'Balkanea-Mobile-Sandbox-Readiness-Update.pptx'
prs = Presentation(PATH)


def shape_by_name(slide, name):
    return [sh for sh in slide.shapes if sh.name == name][0]


def set_single_run_text(shape, new_text):
    tf = shape.text_frame
    first_para = tf.paragraphs[0]
    proto = first_para.runs[0] if first_para.runs else None
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    for r in list(first_para.runs):
        r._r.getparent().remove(r._r)
    run = first_para.add_run()
    run.text = new_text
    if proto is not None:
        run.font.size = proto.font.size
        run.font.bold = proto.font.bold
        run.font.name = proto.font.name
        try:
            if proto.font.color and proto.font.color.type is not None:
                run.font.color.rgb = proto.font.color.rgb
        except Exception:
            pass


def drop_slide_properly(prs, index):
    """Remove a slide and its part/relationship cleanly (avoids partname collisions on re-add)."""
    xml_slides = prs.slides._sldIdLst
    slide_id_el = list(xml_slides)[index]
    rId = slide_id_el.get(qn('r:id'))
    prs.part.drop_rel(rId)
    xml_slides.remove(slide_id_el)


# If this script is being re-run after a previous attempt, drop that slide first
# (identified by its kicker text) rather than leaving an orphan part behind.
for i, s in enumerate(prs.slides):
    for sh in s.shapes:
        if sh.has_text_frame and 'LIVE DEMO' in sh.text_frame.text:
            drop_slide_properly(prs, i)
            break

CREAM = RGBColor(0xFA, 0xF7, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE7, 0xE0, 0xD3)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x65, 0x58)
FOOTER_GRAY = RGBColor(0x9B, 0x93, 0x84)
BROWN = RGBColor(0x7A, 0x3E, 0x0F)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD9)
AMBER_TEXT = RGBColor(0xB8, 0x79, 0x0A)
GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_TEXT = RGBColor(0x2E, 0x6B, 0x57)

layout = prs.slides[3].slide_layout  # same blank layout as the Payments slide
new_slide = prs.slides.add_slide(layout)
new_slide.background.fill.solid()
new_slide.background.fill.fore_color.rgb = CREAM
sp = new_slide.shapes


def textbox(left, top, width, height, text, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = sp.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.name = 'Segoe UI Semibold' if bold else 'Segoe UI'
    run.font.color.rgb = color
    return box


def badge(rect_ref, label, positive=True):
    w = Emu(int(Emu(Pt(5.3)) * len(label) + Emu(Pt(16))))
    h = Inches(0.20)
    left = Emu(int(Emu(rect_ref.left) + Emu(rect_ref.width) - w - Pt(8)))
    top = Emu(int(Emu(rect_ref.top) + Pt(7)))
    b = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    b.adjustments[0] = 0.5
    b.fill.solid()
    b.fill.fore_color.rgb = GREEN_FILL if positive else AMBER_FILL
    b.line.fill.background()
    b.shadow.inherit = False
    tf = b.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(7.5)
    run.font.bold = True
    run.font.name = 'Segoe UI Semibold'
    run.font.color.rgb = GREEN_TEXT if positive else AMBER_TEXT
    return b


def card(left, top, width, height, title, desc, badge_label, badge_positive):
    rect = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.09
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = CARD_BORDER
    rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    pad = Pt(10)
    title_top = Emu(int(top + Inches(0.14)))
    textbox(Emu(int(left + pad)), title_top, Emu(int(width - 2 * pad - Inches(1.6))), Inches(0.22), title, Pt(11.5), bold=True, color=NAVY)
    desc_top = Emu(int(title_top + Inches(0.30)))
    textbox(Emu(int(left + pad)), desc_top, Emu(int(width - 2 * pad)), Emu(int(top + height - desc_top - Pt(8))), desc, Pt(10), color=GRAY)
    if badge_label:
        badge(rect, badge_label, positive=badge_positive)
    return rect


# ---------------------------------------------------------------------------
# Header
# ---------------------------------------------------------------------------
textbox(Inches(0.55), Inches(0.50), Inches(12.23), Inches(0.35), '04  —  LIVE DEMO', Pt(12), bold=True, color=BROWN)
textbox(Inches(0.55), Inches(0.82), Inches(12.23), Inches(0.70), 'Proof the payment gateway itself works', Pt(27), bold=True, color=NAVY)
textbox(Inches(0.55), Inches(1.55), Inches(11.81), Inches(0.55),
        'A live, working demo — with an honest read on what it does and doesn’t tell us about the mobile app.',
        Pt(13.5), color=GRAY)
textbox(Inches(0.55), Inches(7.14), Inches(8.00), Inches(0.30), 'BALKANEA MOBILE  —  SANDBOX / TEST READINESS', Pt(9), color=FOOTER_GRAY)
textbox(Inches(11.78), Inches(7.14), Inches(1.00), Inches(0.30), '?', Pt(9), color=FOOTER_GRAY, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# Cards
# ---------------------------------------------------------------------------
CARD_X, CARD_W = Inches(0.55), Inches(12.23)
TOP0 = Inches(2.35)
CARD_H = Inches(1.02)
GAP = Inches(0.15)

card(CARD_X, TOP0, CARD_W,
    CARD_H,
    'Live demo — try it yourself',
    'https://balkanea-lead-webhook.vercel.app/payment-bridge-demo.html — opens the real Bankart-hosted card form for a live sim-mode order and polls the bridge’s real status endpoint until it confirms. Verified working end-to-end 8/14.',
    'LIVE', True)

card(CARD_X, Emu(int(TOP0 + CARD_H + GAP)), CARD_W, CARD_H,
    'What it proves: the bank side is real',
    'A real WooCommerce order, paid through Bankart’s actual hosted card form (sim mode), with the bridge’s own REST endpoint flipping from pending to success live. This is the one piece of the payment stack that’s fully working today.',
    'CONFIRMED', True)

card(CARD_X, Emu(int(TOP0 + 2 * (CARD_H + GAP))), CARD_W, CARD_H,
    'What it doesn’t prove: the app isn’t wired to it',
    'This runs entirely in a browser against WooCommerce’s guest checkout — the exact WooCommerce-based path that’s been ruled out for the final architecture. The mobile app itself still only runs the simulated flow (activeGateway = simulatedGateway).',
    'CAVEAT', False)

card(CARD_X, Emu(int(TOP0 + 3 * (CARD_H + GAP))), CARD_W, CARD_H,
    'Before you run it live',
    'Each order can only be charged once — a successful test consumes it. Get a fresh pending order and the sim test card details from Hristijan/Ivan ahead of time, not mid-call.',
    'ACTION', False)

prs.save(PATH)
print('Demo slide added. Total slides now:', len(prs.slides))
