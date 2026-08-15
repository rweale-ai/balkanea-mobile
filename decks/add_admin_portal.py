from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PATH = 'Balkanea-Mobile-Sandbox-Readiness-Update.pptx'
prs = Presentation(PATH)
s = prs.slides[6]  # 06 -- ARCHITECTURE
sp = s.shapes

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE7, 0xE0, 0xD3)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x65, 0x58)
FOOTER_GRAY = RGBColor(0x9B, 0x93, 0x84)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD9)
AMBER_TEXT = RGBColor(0xB8, 0x79, 0x0A)
GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_TEXT = RGBColor(0x2E, 0x6B, 0x57)
LINE_GRAY = RGBColor(0xAF, 0xA6, 0x94)


def textbox(left, top, width, height, text, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = sp.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.name = 'Segoe UI Semibold' if bold else 'Segoe UI'
    run.font.color.rgb = color
    return box


def card(left, top, width, height, dashed=False):
    rect = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.09
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = CARD_BORDER
    rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    return rect


def badge(rect_ref, label, positive=True):
    w = Emu(int(Emu(Pt(5.3)) * len(label) + Emu(Pt(16))))
    h = Inches(0.20)
    left = Emu(int(Emu(rect_ref.left) + Emu(rect_ref.width) - w - Pt(8)))
    top = Emu(int(Emu(rect_ref.top) + Pt(7)))
    b = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    b.adjustments[0] = 0.5
    b.fill.solid()
    b.fill.fore_color.rgb = GREEN_FILL if positive else AMBER_FILL
    b.line.fill.background()
    b.shadow.inherit = False
    tf = b.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(7.5)
    run.font.bold = True
    run.font.name = 'Segoe UI Semibold'
    run.font.color.rgb = GREEN_TEXT if positive else AMBER_TEXT
    return b


def component(left, top, width, height, title, desc, badge_label=None, badge_positive=True, dashed=False):
    rect = card(left, top, width, height, dashed=dashed)
    pad = Pt(10)
    title_top = Emu(int(top + (Inches(0.32) if badge_label else Pt(9))))
    textbox(Emu(int(left + pad)), title_top, Emu(int(width - 2 * pad)), Inches(0.22), title, Pt(11.5), bold=True,
            color=(AMBER_TEXT if dashed else NAVY))
    desc_top = Emu(int(title_top + Inches(0.26)))
    textbox(Emu(int(left + pad)), desc_top, Emu(int(width - 2 * pad)), Emu(int(top + height - desc_top - Pt(6))), desc, Pt(9), color=GRAY)
    if badge_label:
        badge(rect, badge_label, positive=badge_positive)
    return rect


def arrow(x1, y1, x2, y2, dashed=False):
    conn = sp.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = AMBER_TEXT if dashed else LINE_GRAY
    conn.line.width = Pt(1.5 if dashed else 1.1)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if dashed:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return conn


# Same layout constants used when this slide was first built.
COL1_X, COL1_W = Inches(0.55), Inches(2.55)
COL2_X, COL2_W = Inches(3.44), Inches(2.55)
BODY_TOP = Inches(2.24)
mobile_h = Inches(1.35)
auth_top = Emu(int(BODY_TOP + mobile_h + Inches(0.24)))
auth_h = Inches(1.05)
backend_top = Inches(2.58)
backend_h = Inches(1.95)

admin_top = Emu(int(auth_top + auth_h + Inches(0.16)))
admin_h = Inches(1.42)

admin = component(COL1_X, admin_top, COL1_W, admin_h,
    'Admin Portal',
    'admin-payments.html — bookings + payment state, Capture/Void stubbed. Doesn’t cover WooCommerce’s refunds/cancellations/invoicing yet — inventorying that gap for Monday.',
    badge_label='PARTIAL', badge_positive=False)

backend_left_mid = (Emu(COL2_X), Emu(int(backend_top + backend_h / 2)))

# Restore the original Mobile App -> Backend connector, accidentally removed during
# cleanup (its left-edge x-coordinate matches the new Admin Portal connector's, since
# both boxes share the Client column's x position).
mobile_right_mid = (Emu(int(COL1_X + COL1_W)), Emu(int(BODY_TOP + mobile_h / 2)))
arrow(mobile_right_mid[0], mobile_right_mid[1], backend_left_mid[0], backend_left_mid[1])

# Admin Portal -> Backend, landing on the same left-mid point the Mobile App arrow uses.
admin_right_mid = (Emu(int(COL1_X + COL1_W)), Emu(int(admin_top + admin_h / 2)))
arrow(admin_right_mid[0], admin_right_mid[1], backend_left_mid[0], backend_left_mid[1])

# Soften the subtitle's "everything else is live or unblocked" claim now that Admin Portal has a known gap.
for sh in s.shapes:
    if sh.has_text_frame and 'One box is still open' in sh.text_frame.text:
        tf = sh.text_frame
        p = tf.paragraphs[0]
        proto = p.runs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        run = p.add_run()
        run.text = ('Two boxes are still open — payments, and the admin portal’s WooCommerce-parity gap. '
                    'Everything else here is either live today or a scoped, unblocked build.')
        run.font.size = proto.font.size
        run.font.bold = proto.font.bold
        run.font.name = proto.font.name
        run.font.color.rgb = proto.font.color.rgb
        break

prs.save(PATH)
print('Saved. Admin Portal box added to Architecture slide.')
