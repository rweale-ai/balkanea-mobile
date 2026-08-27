# -*- coding: utf-8 -*-
"""
2026-08-26 (later) update: layers today's second milestone on top of the
already-current deck (update_content_826.py already fixed the multi-room/
payment/RateHawk-wiring staleness earlier today -- this pass is additive,
not a re-fix of those).

Today's milestone: real amenity filtering, Explore rebuilt to the 3 regions
RateHawk's sandbox actually prices (LA/Paris/Dubai) with real landmark
photos, a backend-wide CORS gap found and fixed (was silently breaking the
whole booking flow on web), and a full real-currency fix (EUR/USD direct
from RateHawk, MKD via peg conversion) threaded through search, room
selection, and payment -- plus search-results-list pricing switched from
simulated to real (one batch RateHawk call per search, not one per hotel).

Also adds a new Performance & Scale slide with real measured RateHawk
sandbox rate limits, and updates the codebase stats to current real commit
counts.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)

# ---------------------------------------------------------------------------
# Part 1: targeted text replacements
# ---------------------------------------------------------------------------
REPLACEMENTS = {
    0: [  # Slide 1 - title subtitle
        (
            "RateHawk sandbox booking and payments are now verified live end-to-end -- three real bookings processed today, including a fixed root-cause payment bug. See Risks for what's still genuinely open. Cost model updated, certification and store submission plan attached.",
            "RateHawk sandbox booking and payments are now verified live end-to-end, including real pricing in the traveler's real selected currency across all three priced regions (LA/Paris/Dubai) -- not just one. See Risks for what's still genuinely open. Cost model updated, certification and store submission plan attached.",
        ),
    ],
    2: [  # Slide 3 - Search & Booking, card 4 description
        (
            "Search \u2192 prebook \u2192 create order \u2192 finish now runs successfully against RateHawk's certification test hotels (Conrad LA, Rosa Bell Motel) -- all 7 mandatory scenarios passed, real order IDs captured. Live-app wiring and the room-by-room guest picker UI are both built and verified -- three real bookings processed through the actual app today.",
            "Search \u2192 prebook \u2192 create order \u2192 finish now runs successfully against RateHawk's certification test hotels (Conrad LA, Rosa Bell Motel) -- all 7 mandatory scenarios passed, real order IDs captured. Live-app wiring and the room-by-room guest picker UI are both built and verified. As of today, real pricing (not simulated) covers both the search results list and room selection across all 3 regions RateHawk's sandbox actually prices, with real amenity filtering and correct multi-currency throughout.",
        ),
    ],
    8: [  # Slide 9 - Architecture: backend + DB boxes
        (
            "Vercel serverless (Node.js). Hotel search proxy, RateHawk mobile workflow (SERP / Prebook / create-order), payment gateway, Salesforce sync.",
            "Vercel serverless (Node.js). Hotel search proxy (real amenity filtering, real multi-currency), RateHawk mobile workflow (SERP / Prebook / create-order), payment gateway, Salesforce sync. CORS gap found and fixed today -- every endpoint had none, silently breaking the whole flow on the web build.",
        ),
        (
            "850,218 hotels / 3.9M rooms, all 10 countries. Region-filtered search, daily refresh from RateHawk dump.",
            "850,218 hotels / 3.9M rooms, all 10 countries. Plus a sibling sandbox schema (same instance): 742 hotels across the 3 regions RateHawk's own sandbox dump actually prices -- real search and room rates, not simulated, as of today.",
        ),
    ],
    10: [  # Slide 11 - What's Been Built
        (
            "\u2022  850K-hotel content DB, region-filtered search",
            "\u2022  850K-hotel content DB + real RateHawk-priced sandbox DB (3 regions), real amenity filtering",
        ),
        (
            "\u2022  Real Bankart/NLB payment integration + voucher/invoice PDF retrieval",
            "\u2022  Real Bankart/NLB payment integration (correct multi-currency: EUR/USD direct from RateHawk, MKD via peg conversion) + voucher/invoice PDF retrieval",
        ),
        (
            "Languages: 2 \u2014 English + Macedonian (Cyrillic), the actual translated UI.\nCurrencies: 9, across 11 countries \u2014 a separate locale/pricing display layer, not UI language. (11 countries is the destination-selector count, not a language count.)",
            "Languages: 2 \u2014 English + Macedonian (Cyrillic), the actual translated UI.\nCurrencies: 9, across 11 countries. For the 3 real-priced regions, EUR/USD come directly from RateHawk (real FX); MKD converts from that real quote via a fixed peg rate.",
        ),
        (
            "ACTUAL CODEBASE, ~10 WEEKS ELAPSED (first commit 6/18 \u2192 8/25)",
            "ACTUAL CODEBASE, ~10 WEEKS ELAPSED (first commit 6/18 \u2192 8/26)",
        ),
        (
            "Mobile: 17 screens \u00b7 12 components \u00b7 27 lib modules \u00b7 ~14,700 lines TS/TSX \u00b7 7 migrations \u00b7 86 commits\nBackend: 12 endpoints \u00b7 12 lib modules \u00b7 ~4,200 lines JS \u00b7 73 commits",
            "Mobile: 17 screens \u00b7 12 components \u00b7 27 lib modules \u00b7 ~14,700 lines TS/TSX \u00b7 7 migrations \u00b7 92 commits\nBackend: 12 endpoints \u00b7 12 lib modules \u00b7 ~4,500 lines JS \u00b7 78 commits",
        ),
    ],
}

changed = 0
for idx, pairs in REPLACEMENTS.items():
    slide = prs.slides[idx]
    for old, new in pairs:
        found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full = ''.join(r.text for r in para.runs)
                if old in full or old.strip() == full.strip():
                    if para.runs:
                        para.runs[0].text = new
                        for r in para.runs[1:]:
                            r.text = ''
                        found = True
                        changed += 1
                        break
            if found:
                break
        if not found:
            print(f"WARNING: text not found on slide index {idx}: {old[:80]}...")

print(f"Replaced {changed} text blocks")

# ---------------------------------------------------------------------------
# Part 2: new Performance & Scale slide, inserted after "11 -- COST"
# (running cost/margin, currently index 11), before Risks (currently 12).
# ---------------------------------------------------------------------------
INSERT_AT = 12


def drop_slide_properly(prs, index):
    xml_slides = prs.slides._sldIdLst
    slide_id_el = list(xml_slides)[index]
    rId = slide_id_el.get(qn('r:id'))
    prs.part.drop_rel(rId)
    xml_slides.remove(slide_id_el)


# Re-run safety
for i, s in enumerate(prs.slides):
    for shp in s.shapes:
        if shp.has_text_frame and "PERFORMANCE & SCALE" in shp.text_frame.text.upper():
            drop_slide_properly(prs, i)
            break

CREAM = RGBColor(0xFA, 0xF7, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE7, 0xE0, 0xD3)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x65, 0x58)
FOOTER_GRAY = RGBColor(0x9B, 0x93, 0x84)
BROWN = RGBColor(0x7A, 0x3E, 0x0F)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD9)
AMBER_TEXT = RGBColor(0xB8, 0x79, 0x0A)
GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_TEXT = RGBColor(0x2E, 0x6B, 0x57)

layout = prs.slides[3].slide_layout
new_slide = prs.slides.add_slide(layout)
new_slide.background.fill.solid()
new_slide.background.fill.fore_color.rgb = CREAM
sp = new_slide.shapes


def textbox(left, top, width, height, text, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT, italic=False):
    box = sp.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = 'Segoe UI Semibold' if bold else 'Segoe UI'
    run.font.color.rgb = color
    return box


def multiline_box(left, top, width, height, lines, size, color=GRAY, line_gap=Pt(4)):
    box = sp.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = line_gap
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.name = 'Segoe UI Semibold' if bold else 'Segoe UI'
        run.font.color.rgb = NAVY if bold else color
    return box


def card(left, top, width, height, title, desc, badge_label, badge_positive):
    rect = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.09
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = CARD_BORDER
    rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    pad = Pt(10)
    title_top = Emu(int(top + Inches(0.12)))
    textbox(Emu(int(left + pad)), title_top, Emu(int(width - 2 * pad - Inches(1.6))), Inches(0.24), title, Pt(12.5), bold=True, color=NAVY)
    desc_top = Emu(int(title_top + Inches(0.30)))
    textbox(Emu(int(left + pad)), desc_top, Emu(int(width - 2 * pad)), Emu(int(top + height - desc_top - Pt(6))), desc, Pt(10), color=GRAY)
    if badge_label:
        w = Emu(int(Emu(Pt(5.3)) * len(badge_label) + Emu(Pt(16))))
        h = Inches(0.20)
        bleft = Emu(int(Emu(left) + Emu(width) - w - Pt(8)))
        btop = Emu(int(Emu(top) + Pt(7)))
        b = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bleft, btop, w, h)
        b.adjustments[0] = 0.5
        b.fill.solid()
        b.fill.fore_color.rgb = GREEN_FILL if badge_positive else AMBER_FILL
        b.line.fill.background()
        b.shadow.inherit = False
        tf = b.text_frame
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = badge_label
        run.font.size = Pt(7.5)
        run.font.bold = True
        run.font.name = 'Segoe UI Semibold'
        run.font.color.rgb = GREEN_TEXT if badge_positive else AMBER_TEXT
    return rect


# Header
textbox(Inches(0.55), Inches(0.50), Inches(12.23), Inches(0.35), '11  \u2014  PERFORMANCE & SCALE', Pt(12), bold=True, color=BROWN)
textbox(Inches(0.55), Inches(0.82), Inches(12.23), Inches(0.70),
        "The real technical ceiling on today's stack, measured live", Pt(27), bold=True, color=NAVY)
textbox(Inches(0.55), Inches(1.55), Inches(11.81), Inches(0.55),
        'RateHawk sandbox rate limits measured directly against the live API today. These are SANDBOX limits, '
        'not confirmed production ones -- see the note below before treating the ceiling as fixed.',
        Pt(13.5), color=GRAY)
textbox(Inches(0.55), Inches(7.14), Inches(8.00), Inches(0.30), 'BALKANEA MOBILE  \u2014  SANDBOX / TEST READINESS', Pt(9), color=FOOTER_GRAY)
textbox(Inches(11.78), Inches(7.14), Inches(1.00), Inches(0.30), '12', Pt(9), color=FOOTER_GRAY, align=PP_ALIGN.RIGHT)

# Three rate-limit stat cards, left side
CARD_Y = Inches(2.35)
CARD_W3 = Inches(3.90)
CARD_H3 = Inches(1.55)
GAP3 = Inches(0.15)

rate_cards = [
    ("search/serp/hotels/", "150 / min", "Search results list -- ONE call prices a whole region (~250 hotels), not one call per hotel.", True),
    ("hotel/prebook/", "30 / min", "Room lock, right before payment. Comfortable headroom above the booking-flow bottleneck.", True),
    ("search/hp/", "10 / min", "Single-hotel room rates -- required on room selection AND before every booking. This is the binding constraint.", False),
]
for i, (endpoint, rate, desc, positive) in enumerate(rate_cards):
    left = Emu(int(Inches(0.55) + i * (CARD_W3 + GAP3)))
    rect = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, CARD_Y, CARD_W3, CARD_H3)
    rect.adjustments[0] = 0.09
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = CARD_BORDER
    rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    pad = Pt(10)
    textbox(Emu(int(left + pad)), Emu(int(CARD_Y + Inches(0.14))), Emu(int(CARD_W3 - 2 * pad)), Inches(0.22), endpoint, Pt(10.5), bold=True, color=BROWN)
    textbox(Emu(int(left + pad)), Emu(int(CARD_Y + Inches(0.42))), Emu(int(CARD_W3 - 2 * pad)), Inches(0.45), rate, Pt(26), bold=True, color=NAVY)
    textbox(Emu(int(left + pad)), Emu(int(CARD_Y + Inches(0.95))), Emu(int(CARD_W3 - 2 * pad)), Emu(int(CARD_H3 - Inches(0.95) - Pt(8))), desc, Pt(9.5), color=GRAY)
    w = Emu(int(Emu(Pt(5.3)) * 18 + Emu(Pt(16))))
    h = Inches(0.20)
    bleft = Emu(int(Emu(left) + Emu(CARD_W3) - w - Pt(8)))
    btop = Emu(int(Emu(CARD_Y) + Pt(7)))
    b = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bleft, btop, w, h)
    b.adjustments[0] = 0.5
    b.fill.solid()
    b.fill.fore_color.rgb = GREEN_FILL if positive else AMBER_FILL
    b.line.fill.background()
    b.shadow.inherit = False
    tf = b.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'BOTTLENECK' if not positive else 'HEADROOM'
    run.font.size = Pt(7.5)
    run.font.bold = True
    run.font.name = 'Segoe UI Semibold'
    run.font.color.rgb = GREEN_TEXT if positive else AMBER_TEXT

# Ceiling callout band
CEILING_Y = Emu(int(CARD_Y + CARD_H3 + Inches(0.20)))
ceiling_rect = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), CEILING_Y, Inches(12.23), Inches(0.95))
ceiling_rect.adjustments[0] = 0.12
ceiling_rect.fill.solid()
ceiling_rect.fill.fore_color.rgb = AMBER_FILL
ceiling_rect.line.fill.background()
ceiling_rect.shadow.inherit = False
textbox(Emu(int(Inches(0.55) + Pt(14))), Emu(int(CEILING_Y + Pt(10))), Inches(11.8), Inches(0.30),
        "10/min on search/hp/ \u2192 600/hour \u2192 ~14,400/day ceiling on real room-rate views, and therefore on bookings, under today's sandbox limits.",
        Pt(14), bold=True, color=BROWN)
textbox(Emu(int(Inches(0.55) + Pt(14))), Emu(int(CEILING_Y + Pt(42))), Inches(11.8), Inches(0.40),
        "This is a technical ceiling, not a target or a forecast \u2014 real booking volume will be governed by guest demand, far below any of these limits.",
        Pt(10.5), color=GRAY, italic=True)

# Bottom row: two supporting-infra cards
BOTTOM_Y = Emu(int(CEILING_Y + Inches(0.95) + Inches(0.18)))
BOTTOM_H = Inches(1.35)
BOTTOM_W = Inches(5.98)
GAP2 = Inches(0.27)

card(Inches(0.55), BOTTOM_Y, BOTTOM_W, BOTTOM_H,
     "Database: comfortable headroom",
     "Supabase Postgres allows 120 concurrent connections; each Vercel function instance currently pools at most 3. "
     "Not a near-term constraint at any traffic level discussed here.",
     'NOT A CONSTRAINT', True)

card(Emu(int(Inches(0.55) + BOTTOM_W + GAP2)), BOTTOM_Y, BOTTOM_W, BOTTOM_H,
     "Infra tier: Vercel Hobby plan",
     "Fine for the current testing phase -- no verified request/day figure exists for it, so none is claimed here. "
     "Real pre-launch item: Hobby has no SLA and stricter usage limits than Pro; upgrade before real production traffic.",
     'PRE-LAUNCH ITEM', False)

# Move the new slide (currently last) into position INSERT_AT
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
new_el = slides[-1]
xml_slides.remove(new_el)
xml_slides.insert(INSERT_AT, new_el)
prs.save(PATH)
print('Inserted "Performance & Scale" slide at index', INSERT_AT, '- total slides now:', len(prs.slides))
