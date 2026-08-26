# -*- coding: utf-8 -*-
from pptx import Presentation

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)
slide = prs.slides[12]

OLD = ("From Ray: Apple Services ID/JWT, Google Play account creation. From Hristijan: Bankart mobile-responsive "
       "fix + billing-address/cardholder-name fields (spec sent 8/25). From Claude: RateHawk app-wiring and "
       "multi-room build are both done and verified -- next is RateHawk's certification questionnaire.")
NEW = ("From Ray: Apple Services ID/JWT, Google Play account. From Hristijan: mobile-responsive fix + "
       "billing/cardholder-name fields (spec sent 8/25). From Claude: RateHawk wiring + multi-room build are "
       "both done -- next is RateHawk's certification questionnaire.")

fixed = 0
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        full = ''.join(r.text for r in para.runs)
        if OLD in full:
            para.runs[0].text = NEW
            for r in para.runs[1:]:
                r.text = ''
            fixed += 1

prs.save(PATH)
print("Fixed:", fixed)
