# -*- coding: utf-8 -*-
"""Fixes the two multi-paragraph replacements update_content_826b.py missed
(Languages/Currencies box and codebase-stats box are 2 separate paragraphs
each, not one \\n-joined paragraph)."""
from pptx import Presentation

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)
s = prs.slides[10]

PARA_REPLACEMENTS = {
    8: {
        1: "Currencies: 9, across 11 countries. For the 3 real-priced regions, EUR/USD come directly from RateHawk (real FX); MKD converts from that real quote via a fixed peg rate.",
    },
    10: {
        0: "Mobile: 17 screens \u00b7 12 components \u00b7 27 lib modules \u00b7 ~14,700 lines TS/TSX \u00b7 7 migrations \u00b7 92 commits",
        1: "Backend: 12 endpoints \u00b7 12 lib modules \u00b7 ~4,500 lines JS \u00b7 78 commits",
    },
}

for shape_idx, para_map in PARA_REPLACEMENTS.items():
    sh = s.shapes[shape_idx]
    for para_idx, new_text in para_map.items():
        p = sh.text_frame.paragraphs[para_idx]
        if p.runs:
            p.runs[0].text = new_text
            for r in p.runs[1:]:
                r.text = ''
            print(f"shape {shape_idx} para {para_idx} -> updated")
        else:
            print(f"WARNING: shape {shape_idx} para {para_idx} has no runs")

prs.save(PATH)
print("Saved.")
