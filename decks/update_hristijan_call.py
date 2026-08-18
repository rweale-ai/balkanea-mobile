from pptx import Presentation

PATH = 'Balkanea-Mobile-Sandbox-Readiness-Update.pptx'
prs = Presentation(PATH)


def replace_text(slide_idx, old_exact, new_text):
    """Find the shape on a slide whose text matches old_exact, replace with
    new_text, collapsing to a single run that reuses the first existing
    run's font (size/bold/name/color) — same pattern already used elsewhere
    in this deck's build scripts."""
    s = prs.slides[slide_idx]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        if sh.text_frame.text.strip() != old_exact.strip():
            continue
        tf = sh.text_frame
        first_para = tf.paragraphs[0]
        proto = first_para.runs[0] if first_para.runs else None
        for p in list(tf.paragraphs[1:]):
            p._p.getparent().remove(p._p)
        for r in list(first_para.runs):
            r._r.getparent().remove(r._r)
        run = first_para.add_run()
        run.text = new_text
        if proto is not None:
            run.font.size = proto.font.size
            run.font.bold = proto.font.bold
            run.font.name = proto.font.name
            try:
                if proto.font.color and proto.font.color.type is not None:
                    run.font.color.rgb = proto.font.color.rgb
            except Exception:
                pass
        return True
    raise RuntimeError(f'Slide {slide_idx}: text not found: {old_exact[:60]!r}')


EDITS = [
    # ---- Slide 2 (index 2) — SEARCH & BOOKING / RateHawk sandbox ----
    (2, 'Hotel database is live; RateHawk sandbox is provisioning',
        'Hotel database is live; sandbox key issued, blocked on IP whitelist'),
    (2, "The hotel-DB half of this is done. What's left: RateHawk's sandbox key, the search-flow decision, and the booking endpoints.",
        "The hotel-DB half of this is done. The sandbox key landed 8/17 — first live calls are failing on IP whitelisting, not the credentials themselves. Search-flow decision and booking endpoints still ahead."),
    (2, 'IN PROGRESS', 'BLOCKED'),
    (2, 'RateHawk sandbox — separate contract confirmed, key incoming',
        'Sandbox key issued (ID 954) — first calls blocked, not by bad credentials'),
    (2, "Ivan's proposed contract slug (balkanea-mobile) was accepted by RateHawk today (8/14) — they're creating the contract now and handing off to their API support team to issue the sandbox key. New wrinkle: RateHawk's API Launch team also requires a certification step (demo + checklist) before issuing a production key — a gate not in the original plan.",
        "Received 8/17: Key ID 954 + access token. Live calls return 401 with the key recognized but the token unmatched — consistent with RateHawk's mandatory IP whitelist, not a bad key. Today's ask for Hristijan: does his server already have a RateHawk-whitelisted IP we can reuse?"),
    (2, 'Search-flow design, and booking endpoints still to build',
        "Search-flow now settled by RateHawk's own rules; booking endpoints + 2 missing test hotels still to build"),
    (2, "Still undecided: web's live-search-then-enrich pattern vs. the DB-first shortcut (query our own DB first, no price shown until a hotel is opened). Backend also needs RateHawk's mobile-specific SERP/Prebook/Webhook workflow, not just web's lib/ratehawk.js pattern with swapped credentials — none of the room-rate, pre-book, or create-order calls exist yet.",
        "Not actually a choice anymore: RateHawk's docs prohibit calling the Content API live during search, so DB-first is required, not just faster. None of the room-rate, pre-book, or create-order calls exist yet. New: certification needs 2 specific test hotels (Conrad LA + one other) that aren't in our DB — confirmed by direct query, a small scoped fix."),

    # ---- Slide 3 (index 3) — PAYMENTS ----
    (3, 'Bank checkout — the WooCommerce path is ruled out',
        'Bank checkout is real and scoped; WooCommerce-parity is the remaining question'),
    (3, 'A working WebView integration got built and verified — then a hard constraint ruled it out. Now waiting on the real alternative.',
        "Hristijan confirmed on 8/14 a direct NLB API integration is real and buildable in 2–3 days — his top priority. The open question narrowed too: not which architecture, but whether WooCommerce's specific capabilities (orders, cancellations, invoicing, status tracking) can be fully recreated outside it."),
    (3, 'NLB-direct — the only path that fits, not yet built',
        "NLB-direct — confirmed real, Hristijan's top priority, 2–3 days to build"),
    (3, "Hristijan flagged an NLB EPOS route that skips WooCommerce order creation entirely — surfaced in an 8/14 meeting with Ivan, but 'needs to be developed.' No doc or timeline exists yet on either side; it may be what 'Balkanea Payment Links' (bpl-nlb-settings) already does, but that page is access-blocked for Hristijan's account.",
        "Confirmed directly by Hristijan on the 8/14 call: a direct API integration works against NLB, skipping WooCommerce entirely — 'two to three days tops' to build, now his top priority ahead of other work. Only remaining blocker: NLB needs to whitelist the server's IP for direct API access."),
    (3, 'This week', 'Today’s call'),
    (3, "Actual NLB documentation, what 'needs to be developed' means concretely, and the key question: does the order-free flow still create a WooCommerce order behind the scenes for reconciliation? That answer decides if the constraint is even satisfiable through this vendor.",
        "Confirm NLB IP-whitelist status/timeline, whether the direct flow still writes a WooCommerce order for reconciliation, and — new — whether Hristijan's server already has a RateHawk-whitelisted IP we could reuse for the mobile sandbox key instead of standing up separate infra."),

    # ---- Slide 9 (index 9) — STILL OPEN ----
    (9, 'Contract slug confirmed (balkanea-mobile, accepted 8/14) — RateHawk is creating the contract now and handing off to their API support team to issue the sandbox key. Next: confirm the key lands, then start the certification checklist.',
        "Key landed 8/17 (ID 954) but first calls return incorrect_credentials — looks like IP whitelisting, not a bad key. Next: get Balkanea's outbound IP whitelisted (checking if Hristijan's existing server IP already is), then start the certification checklist — which also needs 2 test hotels added to our DB first."),
    (9, "Pick one: web's live-search-then-enrich pattern, or the DB-first shortcut — needed before backend work continues.",
        "Resolved by RateHawk's own rules, not a preference: their docs prohibit live Content API calls during search, so DB-first is required. Still needs building."),
    (9, "Not a WebView-to-existing-checkout question anymore — get the real NLB-direct spec from Hristijan/Ivan and confirm it doesn't require a WooCommerce order behind the scenes. Nothing else on payments moves until that answer lands.",
        "Confirmed 8/14: NLB-direct is real, ~2–3 days to build, Hristijan's top priority — blocked only on NLB whitelisting the server IP. Separately: RateHawk's booking flow draws on its own deposit balance, not the guest's card, so RateHawk build work doesn't have to wait on this."),
    (9, 'The real NLB documentation from Hristijan/Ivan',
        'IP-whitelist status for both NLB and RateHawk, from Hristijan'),
]

for slide_idx, old, new in EDITS:
    replace_text(slide_idx, old, new)
    print(f'OK  slide {slide_idx}: {old[:45]!r} -> {new[:45]!r}')

prs.save(PATH)
print('\nSaved.')
