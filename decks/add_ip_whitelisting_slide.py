from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PATH = 'Balkanea-Mobile-Sandbox-Readiness-Update.pptx'
prs = Presentation(PATH)

INSERT_AT = 4  # after 03 -- PAYMENTS, before what is currently 04 -- LIVE DEMO


def drop_slide_properly(prs, index):
    xml_slides = prs.slides._sldIdLst
    slide_id_el = list(xml_slides)[index]
    rId = slide_id_el.get(qn('r:id'))
    prs.part.drop_rel(rId)
    xml_slides.remove(slide_id_el)


# Re-run safety: drop a previously-added copy of this slide first.
for i, s in enumerate(prs.slides):
    for sh in s.shapes:
        if sh.has_text_frame and 'IP WHITELISTING' in sh.text_frame.text:
            drop_slide_properly(prs, i)
            break

CREAM = RGBColor(0xFA, 0xF7, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE7, 0xE0, 0xD3)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x65, 0x58)
FOOTER_GRAY = RGBColor(0x9B, 0x93, 0x84)
BROWN = RGBColor(0x7A, 0x3E, 0x0F)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD9)
AMBER_TEXT = RGBColor(0xB8, 0x79, 0x0A)
GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_TEXT = RGBColor(0x2E, 0x6B, 0x57)

layout = prs.slides[3].slide_layout  # same blank layout used for other content slides
new_slide = prs.slides.add_slide(layout)
new_slide.background.fill.solid()
new_slide.background.fill.fore_color.rgb = CREAM
sp = new_slide.shapes


def textbox(left, top, width, height, text, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = sp.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.name = 'Segoe UI Semibold' if bold else 'Segoe UI'
    run.font.color.rgb = color
    return box


def badge(rect_ref, label, positive=True):
    w = Emu(int(Emu(Pt(5.3)) * len(label) + Emu(Pt(16))))
    h = Inches(0.20)
    left = Emu(int(Emu(rect_ref.left) + Emu(rect_ref.width) - w - Pt(8)))
    top = Emu(int(Emu(rect_ref.top) + Pt(7)))
    b = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    b.adjustments[0] = 0.5
    b.fill.solid()
    b.fill.fore_color.rgb = GREEN_FILL if positive else AMBER_FILL
    b.line.fill.background()
    b.shadow.inherit = False
    tf = b.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(7.5)
    run.font.bold = True
    run.font.name = 'Segoe UI Semibold'
    run.font.color.rgb = GREEN_TEXT if positive else AMBER_TEXT
    return b


def card(left, top, width, height, title, desc, badge_label, badge_positive):
    rect = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.09
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = CARD_BORDER
    rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    pad = Pt(10)
    title_top = Emu(int(top + Inches(0.14)))
    textbox(Emu(int(left + pad)), title_top, Emu(int(width - 2 * pad - Inches(1.6))), Inches(0.22), title, Pt(11.5), bold=True, color=NAVY)
    desc_top = Emu(int(title_top + Inches(0.30)))
    textbox(Emu(int(left + pad)), desc_top, Emu(int(width - 2 * pad)), Emu(int(top + height - desc_top - Pt(8))), desc, Pt(10), color=GRAY)
    if badge_label:
        badge(rect, badge_label, positive=badge_positive)
    return rect


# ---------------------------------------------------------------------------
# Header
# ---------------------------------------------------------------------------
textbox(Inches(0.55), Inches(0.50), Inches(12.23), Inches(0.35), '04  —  IP WHITELISTING', Pt(12), bold=True, color=BROWN)
textbox(Inches(0.55), Inches(0.82), Inches(12.23), Inches(0.70), 'One static egress IP, shared by RateHawk and NLB — or two?', Pt(27), bold=True, color=NAVY)
textbox(Inches(0.55), Inches(1.55), Inches(11.81), Inches(0.55),
        'Both vendors require a whitelisted outbound IP. Neither is confirmed yet — here is what is known, what is not, and where this should live.',
        Pt(13.5), color=GRAY)
textbox(Inches(0.55), Inches(7.14), Inches(8.00), Inches(0.30), 'BALKANEA MOBILE  —  SANDBOX / TEST READINESS', Pt(9), color=FOOTER_GRAY)
textbox(Inches(11.78), Inches(7.14), Inches(1.00), Inches(0.30), '5', Pt(9), color=FOOTER_GRAY, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# Cards
# ---------------------------------------------------------------------------
CARD_X, CARD_W = Inches(0.55), Inches(12.23)
TOP0 = Inches(2.35)
CARD_H = Inches(1.02)
GAP = Inches(0.15)

card(CARD_X, TOP0, CARD_W, CARD_H,
    'Two vendors, one unconfirmed cause',
    "RateHawk's checklist calls IP whitelisting mandatory, but our 401 (incorrect_credentials, api_key_id: 0) reads more like an inactive/unpropagated key than a blocked IP — a real gateway usually returns a distinct ip_not_allowed / 403 for that. Don't build infra on an assumption: email RateHawk support the exact response and ask directly.",
    'UNCONFIRMED', False)

card(CARD_X, Emu(int(TOP0 + CARD_H + GAP)), CARD_W, CARD_H,
    'Recommended: a thin static-IP egress layer, not a bigger server',
    "Decouple the static IP from where the app runs. Keep the backend on Vercel; route only the outbound RateHawk/NLB calls through a small, stateless proxy (self-hosted or a static-IP add-on) that owns the whitelisted IP. One shared IP for both vendors is fine — the proxy carries near-zero load, so splitting it in two buys nothing but a second whitelist request. Reliability comes from running two proxy instances behind one floating IP, not from separate boxes.",
    'RECOMMENDED', True)

card(CARD_X, Emu(int(TOP0 + 2 * (CARD_H + GAP))), CARD_W, CARD_H,
    "Hristijan's test server unblocks this week — it isn't the destination",
    "His 8/18 successful NLB test debit proves *some* origin of his already works. Reusing it is the fastest way to start testing RateHawk's sandbox key too. But it's coupled to the legacy WooCommerce box — fine as a stopgap, wrong as the long-term home given the goal of keeping WooCommerce out of the mobile app's critical path.",
    'STOPGAP', False)

card(CARD_X, Emu(int(TOP0 + 3 * (CARD_H + GAP))), CARD_W, CARD_H,
    'Two answers needed before choosing',
    "From Hristijan: what IP did the 8/18 NLB test call originate from, and does NLB's SIM environment even enforce whitelisting? From RateHawk: is the 401 an IP problem or an inactive key, and what IP should be whitelisted? Both are one message away — architecture choice waits on neither being guessed.",
    'ACTION', False)

prs.save(PATH)

# ---------------------------------------------------------------------------
# Move the new slide (currently last) into position INSERT_AT
# ---------------------------------------------------------------------------
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
new_el = slides[-1]
xml_slides.remove(new_el)
xml_slides.insert(INSERT_AT, new_el)
prs.save(PATH)
print('Inserted new slide at index', INSERT_AT, '- total slides now:', len(prs.slides))
