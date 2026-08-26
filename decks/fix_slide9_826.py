# -*- coding: utf-8 -*-
"""Fixes two issues on the Architecture slide (index 8):
1. balkanea_sbx card text overflows its fixed-height box -- shorten it.
2. The slide's own subtitle still has the stale "treated as live... see
   Risks for what that assumes" framing -- same fix as Slide 1.
"""
from pptx import Presentation

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)
slide = prs.slides[8]

PAIRS = [
    (
        "Bookings, profiles, payment_reference / payment_state (MKD), ratehawk_order_id. Live -- three real bookings written today via the fully-wired RateHawk flow. New: voucher/invoice PDF fetch and automated confirmation emails (guest + business team) via Resend.",
        "Bookings, profiles, payment_reference/payment_state (MKD), ratehawk_order_id. Live -- three real bookings written today. New: voucher/invoice PDF fetch + automated confirmation emails via Resend.",
    ),
    (
        "Payments and RateHawk booking are treated as live for this update (see Risks for what that assumes); the admin portal's WooCommerce-parity gap is the one box still genuinely open below.",
        "Payments and RateHawk booking are now verified live, not just treated as such (see Risks for what's still open); the admin portal's WooCommerce-parity gap is the one box still genuinely open below.",
    ),
]

fixed = 0
for old, new in PAIRS:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = ''.join(r.text for r in para.runs)
            if old in full:
                para.runs[0].text = new
                for r in para.runs[1:]:
                    r.text = ''
                fixed += 1

prs.save(PATH)
print("Fixed:", fixed)
