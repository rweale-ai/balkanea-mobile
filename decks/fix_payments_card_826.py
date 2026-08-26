# -*- coding: utf-8 -*-
"""Shortens the 'Balkanea Payment Bridge' card body on Slide 4 -- the
8/26 replacement text overflowed the card's fixed height and collided
with the card below it."""
from pptx import Presentation

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)
slide = prs.slides[3]

OLD = ("Hristijan's \u2018Balkanea Payment Bridge\u2019 plugin (signed HMAC links, Bankart Gateway V3 hosted card form, "
       "async Notify webhook) is live. Confirmed 8/20 and re-confirmed 8/25 against a real RateHawk hotel, settling "
       "in MKD (Bankart's only supported currency, confirmed directly with Hristijan). 8/25 also found and fixed "
       "the actual root cause of intermittent failures: Hristijan's plugin appends its own timestamp suffix to our "
       "payment reference before submitting to Bankart, silently exceeding Bankart's 50-character limit on "
       "real-hotel bookings. Fixed app-side; charges now clear reliably.")
NEW = ("Hristijan's \u2018Balkanea Payment Bridge\u2019 plugin (signed HMAC links, Bankart Gateway V3 hosted card form, "
       "async Notify webhook) is live. Confirmed 8/20 and re-confirmed 8/25 against a real RateHawk hotel, "
       "settling in MKD (Bankart's only supported currency). 8/25 also found and fixed the root cause of "
       "intermittent failures -- Hristijan's plugin was silently exceeding Bankart's 50-character reference "
       "limit on real-hotel bookings. Fixed app-side; charges now clear reliably.")

fixed = False
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        full = ''.join(r.text for r in para.runs)
        if OLD in full:
            para.runs[0].text = NEW
            for r in para.runs[1:]:
                r.text = ''
            fixed = True

prs.save(PATH)
print("Fixed:", fixed)
