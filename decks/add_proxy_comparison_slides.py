from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PATH = 'Balkanea-Mobile-Sandbox-Readiness-Update.pptx'
prs = Presentation(PATH)

INSERT_AT = 5  # after 04 -- IP WHITELISTING


def drop_slide_properly(prs, index):
    xml_slides = prs.slides._sldIdLst
    slide_id_el = list(xml_slides)[index]
    rId = slide_id_el.get(qn('r:id'))
    prs.part.drop_rel(rId)
    xml_slides.remove(slide_id_el)


# Re-run safety: drop previously-added copies of these two slides first.
for marker in ['PROXY: IVAN', 'ALIGNMENT & RECOMMENDATION']:
    for i, s in enumerate(prs.slides):
        found = False
        for sh in s.shapes:
            if sh.has_text_frame and marker in sh.text_frame.text:
                drop_slide_properly(prs, i)
                found = True
                break
        if found:
            break

CREAM = RGBColor(0xFA, 0xF7, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE7, 0xE0, 0xD3)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x65, 0x58)
FOOTER_GRAY = RGBColor(0x9B, 0x93, 0x84)
BROWN = RGBColor(0x7A, 0x3E, 0x0F)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xD9)
AMBER_TEXT = RGBColor(0xB8, 0x79, 0x0A)
GREEN_FILL = RGBColor(0xE3, 0xEE, 0xE9)
GREEN_TEXT = RGBColor(0x2E, 0x6B, 0x57)
DIVIDER = RGBColor(0xEC, 0xE6, 0xDA)


def new_slide():
    layout = prs.slides[3].slide_layout
    s = prs.slides.add_slide(layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = CREAM
    return s


def textbox(sp, left, top, width, height, text, size, bold=False, color=NAVY, align=PP_ALIGN.LEFT, italic=False):
    box = sp.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = 'Segoe UI Semibold' if bold else 'Segoe UI'
    run.font.color.rgb = color
    return box


def header(sp, num_title, headline, subhead, page_num):
    textbox(sp, Inches(0.55), Inches(0.50), Inches(12.23), Inches(0.35), num_title, Pt(12), bold=True, color=BROWN)
    textbox(sp, Inches(0.55), Inches(0.82), Inches(12.23), Inches(0.62), headline, Pt(25), bold=True, color=NAVY)
    textbox(sp, Inches(0.55), Inches(1.46), Inches(11.81), Inches(0.45), subhead, Pt(12.5), color=GRAY)
    textbox(sp, Inches(0.55), Inches(7.14), Inches(8.00), Inches(0.30), 'BALKANEA MOBILE  —  SANDBOX / TEST READINESS', Pt(9), color=FOOTER_GRAY)
    textbox(sp, Inches(11.78), Inches(7.14), Inches(1.00), Inches(0.30), str(page_num), Pt(9), color=FOOTER_GRAY, align=PP_ALIGN.RIGHT)


def pill(sp, left, top, label, positive=True, neutral=False):
    w = Emu(int(Pt(6.0) * len(label) + Pt(20)))
    h = Inches(0.22)
    b = sp.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    b.adjustments[0] = 0.5
    b.fill.solid()
    b.fill.fore_color.rgb = GREEN_FILL if positive else (RGBColor(0xEE, 0xEC, 0xE6) if neutral else AMBER_FILL)
    b.line.fill.background()
    b.shadow.inherit = False
    tf = b.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.name = 'Segoe UI Semibold'
    run.font.color.rgb = GREEN_TEXT if positive else (GRAY if neutral else AMBER_TEXT)
    return b, w


# =============================================================================
# SLIDE 1: Comparison across the five dimensions
# =============================================================================
s1 = new_slide()
sp1 = s1.shapes

header(sp1, '05  —  PROXY: IVAN VS AWS',
       "Ivan's proxy proposal vs. what we built on AWS",
       'Scoped only to the static-IP requirement for RateHawk — not the broader API-gateway scope in his original doc. Unbiased comparison; both columns are real options.',
       6)

COL_LABEL_W = Inches(2.15)
COL_X = Inches(0.55)
COL1_X = Emu(int(COL_X + COL_LABEL_W))
COL_W = Inches(4.90)
COL2_X = Emu(int(COL1_X + COL_W + Inches(0.18)))
TOP0 = Inches(2.28)
ROW_H = Inches(0.80)
GAP = Inches(0.06)

# Column headers
textbox(sp1, COL1_X, Emu(int(TOP0 - Inches(0.34))), COL_W, Inches(0.28), "IVAN'S PROXY PROPOSAL", Pt(10.5), bold=True, color=BROWN)
textbox(sp1, COL2_X, Emu(int(TOP0 - Inches(0.34))), COL_W, Inches(0.28), "OUR AWS BUILD (LIVE TODAY)", Pt(10.5), bold=True, color=BROWN)

rows = [
    ('COST', 'Not costed in his doc. As specified — VPS + Redis + eventual HA instance — realistically ~$20–50+/mo, plus unspecified engineering time to build it.',
     '~$5–8/mo (single small instance), $0 static IP, $0 TLS. Real, verified spend — not an estimate.', True),
    ('PERFORMANCE', 'Designed for high scale from day one (Redis rate-limiting, BullMQ queueing). Higher ceiling, but unbuilt and untested — solving a scale problem this workload does not have yet.',
     'Single lightweight instance handles this thin forwarding workload with large headroom. RateHawk’s own rate limits are the real ceiling either way.', None),
    ('STABILITY', 'Phase 1 as specified is also single-instance — his own doc defers HA to a future phase. Same posture as ours at this phase.',
     'Single instance, auto-restart on crash, fully rebuildable from source-controlled scripts. Box itself verified working; RateHawk-specific logic not yet built on either side.', None),
    ('SUPPORT', "No hosting provider or support commitment specified. Who builds and operates it long-term was never answered in the proposal — a real, still-open gap.",
     'Real AWS infrastructure (~99.99% published EC2 SLA). MARRA can add paid AWS support later if ever needed. Ownership is resolved: this is ours to run.', True),
    ('SCALABILITY', 'Redis/BullMQ/HA path built in from the start — a genuine advantage once traffic actually needs it.',
     'Same evolution path stays fully open — nothing blocks adding Redis/BullMQ/HA later, exactly when real traffic justifies it. Not paying for it before it is needed.', None),
]

for i, (label, left_text, right_text, aws_wins) in enumerate(rows):
    top = Emu(int(TOP0 + i * (ROW_H + GAP)))
    # row background card spans full width for subtle separation
    rowbg = sp1.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, COL_X, top, Inches(12.23), ROW_H)
    rowbg.adjustments[0] = 0.06
    rowbg.fill.solid()
    rowbg.fill.fore_color.rgb = WHITE
    rowbg.line.color.rgb = CARD_BORDER
    rowbg.line.width = Pt(0.75)
    rowbg.shadow.inherit = False

    textbox(sp1, Emu(int(COL_X + Pt(10))), Emu(int(top + Inches(0.10))), Emu(int(COL_LABEL_W - Pt(16))), Inches(0.5),
            label, Pt(11), bold=True, color=NAVY)

    textbox(sp1, Emu(int(COL1_X + Pt(10))), Emu(int(top + Inches(0.08))), Emu(int(COL_W - Pt(20))), Emu(int(ROW_H - Inches(0.16))),
            left_text, Pt(9.5), color=GRAY)

    textbox(sp1, Emu(int(COL2_X + Pt(10))), Emu(int(top + Inches(0.08))), Emu(int(COL_W - Pt(20))), Emu(int(ROW_H - Inches(0.16))),
            right_text, Pt(9.5), color=(GREEN_TEXT if aws_wins else GRAY))

prs.save(PATH)


# =============================================================================
# SLIDE 2: Analytics alignment, fair points either way, recommendation
# =============================================================================
s2 = new_slide()
sp2 = s2.shapes

header(sp2, '06  —  ALIGNMENT & RECOMMENDATION',
       'Same pattern as Analytics — the fastest path to go-live',
       'Both options are real. Here is the case for each, and why one fits this specific requirement better.',
       7)

CARD_X, CARD_W = Inches(0.55), Inches(12.23)
TOP0b = Inches(2.30)

# Card: Analytics alignment
rect = sp2.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, CARD_X, TOP0b, CARD_W, Inches(1.00))
rect.adjustments[0] = 0.09
rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
rect.line.color.rgb = CARD_BORDER; rect.line.width = Pt(0.75)
rect.shadow.inherit = False
textbox(sp2, Emu(int(CARD_X + Pt(10))), Emu(int(TOP0b + Inches(0.14))), Emu(int(CARD_W - Pt(20) - Inches(1.6))), Inches(0.24),
        'Aligned with the Analytics Architecture, not a one-off', Pt(11.5), bold=True, color=NAVY)
textbox(sp2, Emu(int(CARD_X + Pt(10))), Emu(int(TOP0b + Inches(0.44))), Emu(int(CARD_W - Pt(20))), Inches(0.50),
        "The same static-IP relay pattern — a small, dedicated, isolated instance with a static IP — is already the documented pattern for Alana Analytics' own extraction relay. "
        "Both now live under one AWS Organization, in separate accounts for compliance isolation, following the same internal playbook. Ivan's proposal has no relationship to this "
        "existing, already-proven pattern.", Pt(9.8), color=GRAY)
pill(sp2, Emu(int(CARD_X + CARD_W - Inches(1.55))), Emu(int(TOP0b + Pt(7))), 'MARRA-WIDE', positive=True)

TOP1b = Emu(int(TOP0b + Inches(1.00) + Inches(0.14)))
rect = sp2.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, CARD_X, TOP1b, CARD_W, Inches(1.10))
rect.adjustments[0] = 0.09
rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
rect.line.color.rgb = CARD_BORDER; rect.line.width = Pt(0.75)
rect.shadow.inherit = False
textbox(sp2, Emu(int(CARD_X + Pt(10))), Emu(int(TOP1b + Inches(0.14))), Emu(int(CARD_W - Pt(20) - Inches(1.6))), Inches(0.24),
        "Fair case for Ivan's proposal", Pt(11.5), bold=True, color=NAVY)
textbox(sp2, Emu(int(CARD_X + Pt(10))), Emu(int(TOP1b + Inches(0.44))), Emu(int(CARD_W - Pt(20))), Inches(0.60),
        "If Balkanea's own team builds and runs it, they get full visibility without depending on MARRA — a real advantage for long-term independence. "
        "His design also has a cleaner path to a full API-gateway abstraction later (RateHawk becomes an implementation detail, easier to add other suppliers). "
        "Our AWS build currently lives inside MARRA's Organization — Balkanea's own team doesn't have access today without a deliberate handoff.", Pt(9.8), color=GRAY)
pill(sp2, Emu(int(CARD_X + CARD_W - Inches(1.15))), Emu(int(TOP1b + Pt(7))), 'FAIR POINT', neutral=True)

TOP2b = Emu(int(TOP1b + Inches(1.10) + Inches(0.14)))
rect = sp2.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, CARD_X, TOP2b, CARD_W, Inches(1.55))
rect.adjustments[0] = 0.09
rect.fill.solid(); rect.fill.fore_color.rgb = GREEN_FILL
rect.line.color.rgb = CARD_BORDER; rect.line.width = Pt(0.75)
rect.shadow.inherit = False
textbox(sp2, Emu(int(CARD_X + Pt(10))), Emu(int(TOP2b + Inches(0.14))), Emu(int(CARD_W - Pt(20) - Inches(1.6))), Inches(0.24),
        'Recommendation: use the AWS build to reach go-live', Pt(11.5), bold=True, color=NAVY)
textbox(sp2, Emu(int(CARD_X + Pt(10))), Emu(int(TOP2b + Inches(0.44))), Emu(int(CARD_W - Pt(20))), Inches(1.00),
        "For the actual requirement — a static IP RateHawk will accept — the AWS build is live, verified, and ~$5–8/mo today, with no further engineering or "
        "coordination needed. Ivan's fuller proposal is well-designed for a mature, high-traffic system, but adopting it as specified means more infrastructure, more cost, "
        "and an unresolved question of who builds and runs it — none of which the actual RateHawk traffic requires right now. Fastest, lowest-risk path to go-live: use what's "
        "already built. Revisit Redis/queueing/the full gateway abstraction once real traffic data justifies it — that evolution path stays open either way.", Pt(9.8), color=GRAY)
pill(sp2, Emu(int(CARD_X + CARD_W - Inches(1.85))), Emu(int(TOP2b + Pt(7))), 'RECOMMENDED', positive=True)

prs.save(PATH)

# ---------------------------------------------------------------------------
# Move the two new slides (currently last two) into position INSERT_AT
# ---------------------------------------------------------------------------
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
first_new = slides[-2]
second_new = slides[-1]
xml_slides.remove(first_new)
xml_slides.remove(second_new)
xml_slides.insert(INSERT_AT, first_new)
xml_slides.insert(INSERT_AT + 1, second_new)
prs.save(PATH)
print('Inserted 2 new slides at index', INSERT_AT, '- total slides now:', len(prs.slides))
