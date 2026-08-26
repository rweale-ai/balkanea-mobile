# -*- coding: utf-8 -*-
from pptx import Presentation

PATH = r"C:\Users\raywe\Ray\Balkanea\Mobile\decks\Balkanea-Mobile-Sandbox-Readiness-Update.pptx"
prs = Presentation(PATH)
slide = prs.slides[9]

OLD = ("Estimated $1,000 engineering cost (Claude Code, if metered -- near-zero if on subscription). "
       "Reflects 8/25's additional real-money payment debugging, reliability fixes, and the new "
       "voucher/invoice/email feature; multi-room is built, not excluded.")
NEW = ("Estimated $1,000 engineering cost (Claude Code, if metered -- near-zero on subscription). "
       "Reflects 8/25's added payment debugging, reliability fixes, and the new voucher/invoice/email feature.")

fixed = 0
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        full = ''.join(r.text for r in para.runs)
        if OLD in full:
            para.runs[0].text = NEW
            for r in para.runs[1:]:
                r.text = ''
            fixed += 1

prs.save(PATH)
print("Fixed:", fixed)
